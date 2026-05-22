import os
import re
import math
import json
import zipfile
import tarfile
import logging
import subprocess
import tempfile
from PySide6.QtWidgets import (QWidget, QVBoxLayout, QLabel, QHBoxLayout,
                               QScrollArea, QSizePolicy, QTextEdit, QPushButton,
                               QStackedWidget, QTextBrowser, QProgressBar,
                               QFrame, QApplication, QTableWidget, QTableWidgetItem,
                               QHeaderView, QTabWidget, QSlider)
from PySide6.QtGui import (QFont, QFontMetrics, QIcon, QPixmap, QColor,
                           QTextCharFormat, QTextCursor, QSyntaxHighlighter,
                           QTextDocument, QPainter, QPen, QImage, QMovie)
from PySide6.QtCore import Qt, QSize, QRectF, QTimer, QThread, Signal, Property, QPropertyAnimation, QEasingCurve, QByteArray

from constants import (TEXT_EXTENSIONS, CODE_EXTENSIONS, IMAGE_EXTENSIONS,
                       DOCUMENT_EXTENSIONS, ARCHIVE_EXTENSIONS, VIDEO_EXTENSIONS, AUDIO_EXTENSIONS)
from ..style_constants import COLORS, FONT, RADIUS, BTN, FILE_ICON_MAP
from ..style_manager import (scrollbar_style, label_caption_style, label_micro_style,
                             badge_style, button_small_primary, button_small_secondary)

logger = logging.getLogger(__name__)

PREVIEWABLE_TEXT_EXTS = TEXT_EXTENSIONS | CODE_EXTENSIONS
PREVIEWABLE_IMAGE_EXTS = IMAGE_EXTENSIONS
MAX_TEXT_PREVIEW_SIZE_MB = 2
MAX_PREVIEW_LINES = 500
MAX_DOC_PREVIEW_SIZE_MB = 10
DEFAULT_PDF_PAGES = 5
MAX_DOC_CHARS = 50000
MAX_MEDIA_PREVIEW_SIZE_MB = 200
MAX_IMAGE_PREVIEW_PIXELS = 1920 * 1080
MARKDOWN_EXTS = {'.md'}
HTML_EXTS = {'.html', '.htm'}
EXCEL_RENDER_PAGES = 3
WORD_LARGE_FILE_CONFIRM_MB = 5

AUTO_PREVIEW_EXTS = ARCHIVE_EXTENSIONS
FOLDER_EXT = '__folder__'


class _MarkdownHighlighter(QSyntaxHighlighter):
    KEYWORDS = {'---', '***', '___', '```'}
    HEADING_CHARS = {'#'}

    def highlightBlock(self, text):
        fmt_heading = QTextCharFormat()
        fmt_heading.setForeground(QColor("#c678dd"))
        fmt_heading.setFontWeight(QFont.Weight.Bold)

        fmt_bold = QTextCharFormat()
        fmt_bold.setForeground(QColor("#e5c07b"))
        fmt_bold.setFontWeight(QFont.Weight.Bold)

        fmt_italic = QTextCharFormat()
        fmt_italic.setForeground(QColor("#e5c07b"))
        fmt_italic.setFontItalic(True)

        fmt_code = QTextCharFormat()
        fmt_code.setForeground(QColor("#98c379"))
        fmt_code.setBackground(QColor(COLORS.BG_TERTIARY))

        fmt_link = QTextCharFormat()
        fmt_link.setForeground(QColor("#61afef"))

        fmt_list = QTextCharFormat()
        fmt_list.setForeground(QColor("#d19a66"))

        fmt_hr = QTextCharFormat()
        fmt_hr.setForeground(QColor("#5c6370"))

        stripped = text.lstrip()
        if stripped.startswith('#'):
            level = 0
            for ch in stripped:
                if ch == '#':
                    level += 1
                else:
                    break
            if 1 <= level <= 6 and level < len(stripped) and stripped[level] == ' ':
                self.setFormat(0, len(text), fmt_heading)
                return

        if stripped.startswith('```'):
            self.setFormat(0, len(text), fmt_code)
            return

        if stripped in ('---', '***', '___'):
            self.setFormat(0, len(text), fmt_hr)
            return

        if stripped.startswith(('- ', '* ', '+ ')):
            self.setFormat(0, len(stripped) - len(stripped.lstrip('-*+ ')) + 1, fmt_list)

        if len(stripped) > 1 and stripped[0].isdigit() and '. ' in stripped[:5]:
            dot_pos = stripped.index('. ')
            self.setFormat(0, dot_pos + 1, fmt_list)

        i = 0
        while i < len(text):
            if text[i] == '`' and i + 1 < len(text):
                j = text.find('`', i + 1)
                if j != -1:
                    self.setFormat(i, j - i + 1, fmt_code)
                    i = j + 1
                    continue
            elif text[i] == '[':
                link_end = text.find(')', i)
                if link_end != -1 and '](' in text[i:link_end]:
                    self.setFormat(i, link_end - i + 1, fmt_link)
                    i = link_end + 1
                    continue
            elif text[i:i+2] == '**':
                end = text.find('**', i + 2)
                if end != -1:
                    self.setFormat(i, end - i + 2, fmt_bold)
                    i = end + 2
                    continue
            elif text[i:i+2] == '__':
                end = text.find('__', i + 2)
                if end != -1:
                    self.setFormat(i, end - i + 2, fmt_bold)
                    i = end + 2
                    continue
            elif text[i] == '*' and i + 1 < len(text) and text[i+1] != '*':
                end = text.find('*', i + 1)
                if end != -1:
                    self.setFormat(i, end - i + 1, fmt_italic)
                    i = end + 1
                    continue
            elif text[i] == '_' and i + 1 < len(text) and text[i+1] != '_':
                end = text.find('_', i + 1)
                if end != -1:
                    self.setFormat(i, end - i + 1, fmt_italic)
                    i = end + 1
                    continue
            i += 1


class _PythonHighlighter(QSyntaxHighlighter):
    KEYWORDS = {'False', 'None', 'True', 'and', 'as', 'assert', 'async', 'await',
                'break', 'class', 'continue', 'def', 'del', 'elif', 'else', 'except',
                'finally', 'for', 'from', 'global', 'if', 'import', 'in', 'is',
                'lambda', 'nonlocal', 'not', 'or', 'pass', 'raise', 'return', 'try',
                'while', 'with', 'yield'}
    BUILTINS = {'print', 'len', 'range', 'int', 'str', 'list', 'dict', 'set', 'tuple',
                'float', 'bool', 'type', 'isinstance', 'hasattr', 'getattr', 'setattr',
                'super', 'self', 'cls', 'open', 'input', 'abs', 'max', 'min', 'sum',
                'enumerate', 'zip', 'map', 'filter', 'sorted', 'reversed', 'any', 'all'}

    def highlightBlock(self, text):
        fmt_keyword = QTextCharFormat()
        fmt_keyword.setForeground(QColor("#c678dd"))
        fmt_keyword.setFontWeight(QFont.Weight.Bold)

        fmt_builtin = QTextCharFormat()
        fmt_builtin.setForeground(QColor("#e5c07b"))

        fmt_string = QTextCharFormat()
        fmt_string.setForeground(QColor("#98c379"))

        fmt_comment = QTextCharFormat()
        fmt_comment.setForeground(QColor("#5c6370"))
        fmt_comment.setFontItalic(True)

        fmt_number = QTextCharFormat()
        fmt_number.setForeground(QColor("#d19a66"))

        fmt_decorator = QTextCharFormat()
        fmt_decorator.setForeground(QColor("#e5c07b"))
        fmt_decorator.setFontItalic(True)

        i = 0
        while i < len(text):
            if text[i] == '#':
                self.setFormat(i, len(text) - i, fmt_comment)
                break
            elif text[i] in '"\'':
                quote = text[i]
                j = i + 1
                while j < len(text):
                    if text[j] == '\\':
                        j += 2
                        continue
                    if text[j] == quote:
                        j += 1
                        break
                    j += 1
                self.setFormat(i, j - i, fmt_string)
                i = j
            elif text[i].isalpha() or text[i] == '_':
                j = i
                while j < len(text) and (text[j].isalnum() or text[j] == '_'):
                    j += 1
                word = text[i:j]
                if word in self.KEYWORDS:
                    self.setFormat(i, j - i, fmt_keyword)
                elif word in self.BUILTINS:
                    self.setFormat(i, j - i, fmt_builtin)
                i = j
            elif text[i] == '@':
                j = i + 1
                while j < len(text) and (text[j].isalnum() or text[j] == '_'):
                    j += 1
                self.setFormat(i, j - i, fmt_decorator)
                i = j
            elif text[i].isdigit() or (text[i] == '.' and i + 1 < len(text) and text[i + 1].isdigit()):
                j = i
                has_dot = False
                while j < len(text) and (text[j].isdigit() or (text[j] == '.' and not has_dot)):
                    if text[j] == '.':
                        has_dot = True
                    j += 1
                self.setFormat(i, j - i, fmt_number)
                i = j
            else:
                i += 1


class _JSHighlighter(QSyntaxHighlighter):
    KEYWORDS = {'break', 'case', 'catch', 'class', 'const', 'continue', 'debugger',
                'default', 'delete', 'do', 'else', 'export', 'extends', 'finally',
                'for', 'function', 'if', 'import', 'in', 'instanceof', 'let', 'new',
                'of', 'return', 'super', 'switch', 'this', 'throw', 'try', 'typeof',
                'var', 'void', 'while', 'with', 'yield', 'async', 'await', 'from',
                'as', 'static', 'get', 'set'}
    BUILTINS = {'console', 'document', 'window', 'Math', 'JSON', 'Promise', 'Array',
                'Object', 'String', 'Number', 'Boolean', 'Map', 'Set', 'Date', 'Error',
                'true', 'false', 'null', 'undefined', 'NaN', 'Infinity'}

    def highlightBlock(self, text):
        fmt_keyword = QTextCharFormat()
        fmt_keyword.setForeground(QColor("#c678dd"))
        fmt_keyword.setFontWeight(QFont.Weight.Bold)
        fmt_string = QTextCharFormat()
        fmt_string.setForeground(QColor("#98c379"))
        fmt_comment = QTextCharFormat()
        fmt_comment.setForeground(QColor("#5c6370"))
        fmt_comment.setFontItalic(True)
        fmt_number = QTextCharFormat()
        fmt_number.setForeground(QColor("#d19a66"))

        i = 0
        while i < len(text):
            if text[i:i+2] == '//':
                self.setFormat(i, len(text) - i, fmt_comment)
                break
            elif text[i] in '"\'`':
                quote = text[i]
                j = i + 1
                while j < len(text):
                    if text[j] == '\\':
                        j += 2
                        continue
                    if text[j] == quote:
                        j += 1
                        break
                    j += 1
                self.setFormat(i, j - i, fmt_string)
                i = j
            elif text[i].isalpha() or text[i] == '_':
                j = i
                while j < len(text) and (text[j].isalnum() or text[j] == '_'):
                    j += 1
                word = text[i:j]
                if word in self.KEYWORDS:
                    self.setFormat(i, j - i, fmt_keyword)
                elif word in self.BUILTINS:
                    fmt_b = QTextCharFormat()
                    fmt_b.setForeground(QColor("#e5c07b"))
                    self.setFormat(i, j - i, fmt_b)
                i = j
            elif text[i].isdigit():
                j = i
                while j < len(text) and (text[j].isdigit() or text[j] == '.'):
                    j += 1
                self.setFormat(i, j - i, fmt_number)
                i = j
            else:
                i += 1


class _GenericCodeHighlighter(QSyntaxHighlighter):
    KEYWORDS_MAP = {
        '.java': {'abstract', 'assert', 'boolean', 'break', 'byte', 'case', 'catch',
                  'char', 'class', 'continue', 'default', 'do', 'double', 'else',
                  'enum', 'extends', 'final', 'finally', 'float', 'for', 'if',
                  'implements', 'import', 'instanceof', 'int', 'interface', 'long',
                  'native', 'new', 'package', 'private', 'protected', 'public',
                  'return', 'short', 'static', 'strictfp', 'super', 'switch',
                  'synchronized', 'this', 'throw', 'throws', 'transient', 'try',
                  'void', 'volatile', 'while', 'true', 'false', 'null'},
        '.c': {'auto', 'break', 'case', 'char', 'const', 'continue', 'default', 'do',
               'double', 'else', 'enum', 'extern', 'float', 'for', 'goto', 'if',
               'int', 'long', 'register', 'return', 'short', 'signed', 'sizeof',
               'static', 'struct', 'switch', 'typedef', 'union', 'unsigned', 'void',
               'volatile', 'while'},
        '.cpp': {'auto', 'break', 'case', 'char', 'class', 'const', 'continue',
                 'default', 'delete', 'do', 'double', 'else', 'enum', 'explicit',
                 'extern', 'float', 'for', 'friend', 'goto', 'if', 'inline', 'int',
                 'long', 'mutable', 'namespace', 'new', 'operator', 'private',
                 'protected', 'public', 'register', 'return', 'short', 'signed',
                 'sizeof', 'static', 'struct', 'switch', 'template', 'this', 'throw',
                 'try', 'typedef', 'typename', 'union', 'unsigned', 'using',
                 'virtual', 'void', 'volatile', 'while', 'nullptr', 'true', 'false'},
        '.h': {'auto', 'break', 'case', 'char', 'class', 'const', 'continue',
               'default', 'delete', 'do', 'double', 'else', 'enum', 'explicit',
               'extern', 'float', 'for', 'friend', 'goto', 'if', 'inline', 'int',
               'long', 'mutable', 'namespace', 'new', 'operator', 'private',
               'protected', 'public', 'register', 'return', 'short', 'signed',
               'sizeof', 'static', 'struct', 'switch', 'template', 'this', 'throw',
               'try', 'typedef', 'typename', 'union', 'unsigned', 'using',
               'virtual', 'void', 'volatile', 'while', 'nullptr', 'true', 'false'},
        '.go': {'break', 'case', 'chan', 'const', 'continue', 'default', 'defer',
                'else', 'fallthrough', 'for', 'func', 'go', 'goto', 'if', 'import',
                'interface', 'map', 'package', 'range', 'return', 'select', 'struct',
                'switch', 'type', 'var', 'true', 'false', 'nil'},
        '.rs': {'as', 'async', 'await', 'break', 'const', 'continue', 'crate', 'dyn',
                'else', 'enum', 'extern', 'fn', 'for', 'if', 'impl', 'in', 'let',
                'loop', 'match', 'mod', 'move', 'mut', 'pub', 'ref', 'return', 'self',
                'Self', 'static', 'struct', 'super', 'trait', 'type', 'unsafe', 'use',
                'where', 'while', 'true', 'false'},
        '.rb': {'BEGIN', 'END', 'alias', 'and', 'begin', 'break', 'case', 'class',
                'def', 'defined?', 'do', 'else', 'elsif', 'end', 'ensure', 'false',
                'for', 'if', 'in', 'module', 'next', 'nil', 'not', 'or', 'redo',
                'rescue', 'retry', 'return', 'self', 'super', 'then', 'true',
                'undef', 'unless', 'until', 'when', 'while', 'yield'},
        '.php': {'abstract', 'and', 'array', 'as', 'break', 'callable', 'case',
                 'catch', 'class', 'clone', 'const', 'continue', 'declare', 'default',
                 'die', 'do', 'echo', 'else', 'elseif', 'empty', 'enddeclare',
                 'endfor', 'endforeach', 'endif', 'endswitch', 'endwhile', 'eval',
                 'exit', 'extends', 'final', 'finally', 'fn', 'for', 'foreach',
                 'function', 'global', 'goto', 'if', 'implements', 'include',
                 'instanceof', 'interface', 'isset', 'list', 'namespace', 'new',
                 'or', 'print', 'private', 'protected', 'public', 'require',
                 'return', 'static', 'switch', 'throw', 'trait', 'try', 'unset',
                 'use', 'var', 'while', 'xor', 'yield', 'true', 'false', 'null'},
        '.sql': {'SELECT', 'FROM', 'WHERE', 'INSERT', 'INTO', 'VALUES', 'UPDATE',
                 'SET', 'DELETE', 'CREATE', 'TABLE', 'ALTER', 'DROP', 'INDEX',
                 'JOIN', 'INNER', 'LEFT', 'RIGHT', 'OUTER', 'ON', 'AND', 'OR',
                 'NOT', 'IN', 'EXISTS', 'BETWEEN', 'LIKE', 'ORDER', 'BY', 'GROUP',
                 'HAVING', 'LIMIT', 'OFFSET', 'UNION', 'ALL', 'AS', 'DISTINCT',
                 'COUNT', 'SUM', 'AVG', 'MIN', 'MAX', 'NULL', 'IS', 'PRIMARY',
                 'KEY', 'FOREIGN', 'REFERENCES', 'CONSTRAINT', 'DEFAULT', 'CHECK',
                 'UNIQUE', 'AUTO_INCREMENT', 'IF', 'THEN', 'ELSE', 'END', 'CASE',
                 'WHEN', 'VIEW', 'TRIGGER', 'PROCEDURE', 'FUNCTION', 'BEGIN',
                 'COMMIT', 'ROLLBACK', 'GRANT', 'REVOKE'},
        '.sh': {'if', 'then', 'else', 'elif', 'fi', 'case', 'esac', 'for', 'while',
                'until', 'do', 'done', 'in', 'function', 'select', 'time', 'coproc',
                'return', 'exit', 'break', 'continue', 'declare', 'export', 'local',
                'readonly', 'typeset', 'unset', 'true', 'false'},
        '.bat': {'echo', 'set', 'if', 'else', 'for', 'do', 'goto', 'call', 'exit',
                 'pause', 'rem', 'start', 'title', 'color', 'cls', 'copy', 'del',
                 'move', 'ren', 'mkdir', 'rmdir', 'cd', 'dir', 'type', 'find',
                 'findstr', 'sort', 'more', 'choice', 'timeout', 'taskkill',
                 'not', 'exist', 'defined', 'equ', 'neq', 'lss', 'leq', 'gtr', 'geq'},
        '.ps1': {'if', 'else', 'elseif', 'switch', 'for', 'foreach', 'while', 'do',
                 'until', 'break', 'continue', 'return', 'throw', 'try', 'catch',
                 'finally', 'function', 'filter', 'workflow', 'class', 'enum',
                 'param', 'dynamicparam', 'begin', 'process', 'end', 'in',
                 'from', 'where', 'select', 'sort', 'group', 'measure', 'foreach-object',
                 'where-object', 'true', 'false', '$null'},
        '.html': {'html', 'head', 'body', 'div', 'span', 'p', 'a', 'img', 'ul', 'ol',
                  'li', 'table', 'tr', 'td', 'th', 'form', 'input', 'button',
                  'select', 'option', 'textarea', 'script', 'style', 'link', 'meta',
                  'title', 'header', 'footer', 'nav', 'section', 'article', 'aside',
                  'main', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'br', 'hr', 'strong',
                  'em', 'code', 'pre', 'blockquote'},
        '.css': {'color', 'background', 'margin', 'padding', 'border', 'font',
                 'display', 'position', 'width', 'height', 'top', 'left', 'right',
                 'bottom', 'flex', 'grid', 'align', 'justify', 'overflow', 'opacity',
                 'transition', 'transform', 'animation', 'box-shadow', 'text-align',
                 'line-height', 'font-size', 'font-weight', 'font-family', 'z-index',
                 'important', 'none', 'auto', 'inherit', 'initial', 'unset', 'relative',
                 'absolute', 'fixed', 'sticky', 'block', 'inline', 'flex', 'grid',
                 'hidden', 'visible', 'scroll', 'center', 'space-between', 'column',
                 'row', 'wrap', 'nowrap'},
        '.json': {'true', 'false', 'null'},
    }

    def __init__(self, doc, ext):
        self._ext = ext
        super().__init__(doc)

    def highlightBlock(self, text):
        keywords = self.KEYWORDS_MAP.get(self._ext, set())
        if not keywords:
            return
        fmt_keyword = QTextCharFormat()
        fmt_keyword.setForeground(QColor("#c678dd"))
        fmt_keyword.setFontWeight(QFont.Weight.Bold)
        fmt_string = QTextCharFormat()
        fmt_string.setForeground(QColor("#98c379"))
        fmt_comment = QTextCharFormat()
        fmt_comment.setForeground(QColor("#5c6370"))
        fmt_comment.setFontItalic(True)
        fmt_number = QTextCharFormat()
        fmt_number.setForeground(QColor("#d19a66"))

        i = 0
        while i < len(text):
            if self._ext in ('.sh', '.bat', '.ps1') and text[i] == '#':
                self.setFormat(i, len(text) - i, fmt_comment)
                break
            elif self._ext in ('.sql',) and text[i:i+2] == '--':
                self.setFormat(i, len(text) - i, fmt_comment)
                break
            elif self._ext in ('.html',) and '<!--' in text[i:]:
                idx = text.index('<!--', i)
                end = text.find('-->', idx)
                if end == -1:
                    self.setFormat(idx, len(text) - idx, fmt_comment)
                else:
                    self.setFormat(idx, end + 3 - idx, fmt_comment)
                break
            elif self._ext in ('.css',) and text[i:i+2] == '/*':
                end = text.find('*/', i + 2)
                if end == -1:
                    self.setFormat(i, len(text) - i, fmt_comment)
                else:
                    self.setFormat(i, end + 2 - i, fmt_comment)
                i = end + 2 if end != -1 else len(text)
                continue
            elif text[i] in '"\'':
                quote = text[i]
                j = i + 1
                while j < len(text):
                    if text[j] == '\\':
                        j += 2
                        continue
                    if text[j] == quote:
                        j += 1
                        break
                    j += 1
                self.setFormat(i, j - i, fmt_string)
                i = j
            elif text[i].isalpha() or text[i] == '_':
                j = i
                while j < len(text) and (text[j].isalnum() or text[j] in '_-'):
                    j += 1
                word = text[i:j]
                if self._ext in ('.sql',) and word.upper() in {k.upper() for k in keywords}:
                    self.setFormat(i, j - i, fmt_keyword)
                elif word in keywords:
                    self.setFormat(i, j - i, fmt_keyword)
                i = j
            elif text[i].isdigit():
                j = i
                while j < len(text) and (text[j].isdigit() or text[j] in '.%pxem'):
                    j += 1
                self.setFormat(i, j - i, fmt_number)
                i = j
            else:
                i += 1


class _JSONHighlighter(QSyntaxHighlighter):
    """JSON 语法高亮器，区分键名和字符串值"""

    def highlightBlock(self, text):
        fmt_key = QTextCharFormat()
        fmt_key.setForeground(QColor("#e5c07b"))
        fmt_key.setFontWeight(QFont.Weight.Bold)

        fmt_string = QTextCharFormat()
        fmt_string.setForeground(QColor("#98c379"))

        fmt_number = QTextCharFormat()
        fmt_number.setForeground(QColor("#d19a66"))

        fmt_bool = QTextCharFormat()
        fmt_bool.setForeground(QColor("#c678dd"))
        fmt_bool.setFontWeight(QFont.Weight.Bold)

        fmt_null = QTextCharFormat()
        fmt_null.setForeground(QColor("#c678dd"))
        fmt_null.setFontItalic(True)

        i = 0
        while i < len(text):
            if text[i] in '"\'':
                quote = text[i]
                j = i + 1
                while j < len(text):
                    if text[j] == '\\':
                        j += 2
                        continue
                    if text[j] == quote:
                        j += 1
                        break
                    j += 1
                # 判断是否为键名（后面紧跟冒号）
                k = j
                while k < len(text) and text[k] in ' \t':
                    k += 1
                if k < len(text) and text[k] == ':':
                    self.setFormat(i, j - i, fmt_key)
                else:
                    self.setFormat(i, j - i, fmt_string)
                i = j
            elif text[i].isdigit() or text[i] == '-':
                j = i + 1
                has_dot = False
                has_e = False
                while j < len(text):
                    if text[j].isdigit():
                        j += 1
                    elif text[j] == '.' and not has_dot:
                        has_dot = True
                        j += 1
                    elif text[j] in 'eE' and not has_e:
                        has_e = True
                        j += 1
                        if j < len(text) and text[j] in '+-':
                            j += 1
                    else:
                        break
                self.setFormat(i, j - i, fmt_number)
                i = j
            elif text[i].isalpha() or text[i] == '_':
                j = i
                while j < len(text) and (text[j].isalnum() or text[j] == '_'):
                    j += 1
                word = text[i:j]
                if word in ('true', 'false'):
                    self.setFormat(i, j - i, fmt_bool)
                elif word == 'null':
                    self.setFormat(i, j - i, fmt_null)
                i = j
            else:
                i += 1


def _create_highlighter(doc, ext):
    if ext in ('.py',):
        return _PythonHighlighter(doc)
    elif ext in ('.js', '.ts'):
        return _JSHighlighter(doc)
    elif ext == '.md':
        return _MarkdownHighlighter(doc)
    elif ext == '.json':
        return _JSONHighlighter(doc)
    elif ext in _GenericCodeHighlighter.KEYWORDS_MAP:
        return _GenericCodeHighlighter(doc, ext)
    return None


class ElidedFileNameLabel(QLabel):
    def __init__(self, text="", parent=None):
        super().__init__(text, parent)
        self._full_text = text

    def setText(self, text):
        self._full_text = text
        super().setText(text)

    def resizeEvent(self, event):
        super().resizeEvent(event)
        self._update_elided()

    def _update_elided(self):
        fm = self.fontMetrics()
        available = self.width()
        elided = fm.elidedText(self._full_text, Qt.TextElideMode.ElideMiddle, available)
        super().setText(elided)

    def minimumSizeHint(self):
        return QSize(0, super().minimumSizeHint().height())

    def sizeHint(self):
        return QSize(0, super().minimumSizeHint().height())


class _FolderEntryRow(QWidget):
    def __init__(self, name: str, is_dir: bool, parent=None):
        super().__init__(parent)
        layout = QHBoxLayout(self)
        layout.setContentsMargins(6, 3, 6, 3)
        layout.setSpacing(6)

        icon_label = QLabel()
        icon_label.setFixedSize(16, 16)
        icon_label.setStyleSheet("border: none; background: transparent;")
        if is_dir:
            icon_label.setPixmap(QIcon("icons/doctype/Folder.svg").pixmap(QSize(16, 16)))
        else:
            ext = os.path.splitext(name)[1].lower()
            icon_name = FILE_ICON_MAP.get(ext, 'doctype/File.svg')
            icon_label.setPixmap(QIcon(f"icons/{icon_name}").pixmap(QSize(16, 16)))

        name_label = QLabel(name)
        name_label.setStyleSheet(f"""
            color: {COLORS.TEXT_SECONDARY};
            font-size: {FONT.MICRO_PT}px;
            border: none;
            background: transparent;
        """)

        layout.addWidget(icon_label)
        layout.addWidget(name_label, 1)
        self.setStyleSheet("QWidget { background: transparent; }")


class _PreviewLoadingSpinner(QWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        self._angle = 0
        self._timer = QTimer(self)
        self._timer.setInterval(40)
        self._timer.timeout.connect(self._rotate)
        self._dot_count = 8
        self._dot_radius = 3
        self._radius = 14
        self.setFixedSize(48, 48)
        self.setVisible(False)

    def start(self):
        self._angle = 0
        self.setVisible(True)
        self._timer.start()

    def stop(self):
        self._timer.stop()
        self.setVisible(False)

    def _rotate(self):
        self._angle = (self._angle + 45) % 360
        self.update()

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.setRenderHint(QPainter.RenderHint.Antialiasing)
        center_x = self.width() / 2
        center_y = self.height() / 2
        for i in range(self._dot_count):
            angle = (self._angle + i * (360 / self._dot_count)) % 360
            alpha = int(255 * (1 - i / self._dot_count))
            rad = angle * math.pi / 180
            x = center_x + self._radius * math.cos(rad)
            y = center_y + self._radius * math.sin(rad)
            color = QColor(COLORS.BRAND)
            color.setAlpha(alpha)
            painter.setPen(Qt.PenStyle.NoPen)
            painter.setBrush(color)
            painter.drawEllipse(int(x - self._dot_radius), int(y - self._dot_radius),
                              self._dot_radius * 2, self._dot_radius * 2)
        painter.end()


class _PdfRenderWorker(QThread):
    """PDF 渲染工作线程，一次性以高 DPI 渲染页面为 QPixmap 数据"""
    page_rendered = Signal(int, QByteArray, int, int)  # page_idx, png_data, width, height
    render_finished = Signal(int)

    def __init__(self, file_path, max_pages, start_page=0, dpi=150, parent=None):
        super().__init__(parent)
        self._file_path = file_path
        self._max_pages = max_pages
        self._start_page = start_page
        self._dpi = dpi
        self._cancelled = False

    def cancel(self):
        self._cancelled = True

    def run(self):
        try:
            import fitz
            doc = fitz.open(self._file_path)
            total = len(doc)
            end_page = min(self._start_page + self._max_pages, total)
            zoom = self._dpi / 72.0
            mat = fitz.Matrix(zoom, zoom)
            for i in range(self._start_page, end_page):
                if self._cancelled:
                    break
                page = doc[i]
                pix = page.get_pixmap(matrix=mat)
                png_data = pix.tobytes("png")
                ba = QByteArray(png_data)
                self.page_rendered.emit(i, ba, pix.width, pix.height)
            doc.close()
            self.render_finished.emit(total)
        except Exception as e:
            logger.warning(f"PDF render worker error: {e}")
            self.render_finished.emit(0)


class _ExcelDataWorker(QThread):
    """读取 Excel 数据的工作线程，返回结构化数据供 QTableWidget 使用"""
    sheet_data_ready = Signal(int, list, list, str)  # idx, headers, rows, sheet_name
    all_sheets_done = Signal(int)  # total sheet count

    def __init__(self, file_path, max_sheets=EXCEL_RENDER_PAGES, parent=None):
        super().__init__(parent)
        self._file_path = file_path
        self._max_sheets = max_sheets
        self._cancelled = False

    def cancel(self):
        self._cancelled = True

    def run(self):
        try:
            from openpyxl import load_workbook
            from openpyxl.utils import get_column_letter
            wb = load_workbook(self._file_path, read_only=True, data_only=True)
            total = len(wb.sheetnames)
            for idx, sheet_name in enumerate(wb.sheetnames):
                if self._cancelled or idx >= self._max_sheets:
                    break
                ws = wb[sheet_name]
                headers = []
                rows = []
                max_cols = 0
                row_count = 0
                for row in ws.iter_rows(values_only=True):
                    if self._cancelled:
                        break
                    cells = [str(c) if c is not None else '' for c in row]
                    rows.append(cells)
                    max_cols = max(max_cols, len(cells))
                    row_count += 1
                    if row_count > 200:
                        break
                if not rows or self._cancelled:
                    continue
                max_cols = min(max_cols, 26)
                headers = [get_column_letter(c + 1) for c in range(max_cols)]
                # 截断每行到 max_cols 列
                for i in range(len(rows)):
                    row = rows[i]
                    if len(row) < max_cols:
                        row.extend([''] * (max_cols - len(row)))
                    else:
                        rows[i] = row[:max_cols]
                self.sheet_data_ready.emit(idx, headers, rows, sheet_name)
            wb.close()
            self.all_sheets_done.emit(total)
        except Exception as e:
            logger.warning(f"Excel data worker error: {e}")
            self.all_sheets_done.emit(0)


class _PreviewModeSlider(QWidget):
    """预览模式二段式分段滑块，类似搜索栏的分段选择器"""
    mode_changed = Signal(int)

    def __init__(self, labels: list, parent=None):
        super().__init__(parent)
        if len(labels) != 2:
            raise ValueError("_PreviewModeSlider 只支持二段式切换")
        self._labels = labels
        self._current_index = 0
        self._slider_pos = 0.0
        self._hovered_segment = -1
        self.setFixedHeight(24)
        self._slider_anim = QPropertyAnimation(self, b"slider_pos")
        self._slider_anim.setDuration(200)
        self._slider_anim.setEasingCurve(QEasingCurve.Type.OutCubic)
        self.setCursor(Qt.CursorShape.PointingHandCursor)
        self.setSizePolicy(QSizePolicy.Policy.Fixed, QSizePolicy.Policy.Fixed)
        self.setMouseTracking(True)

    def get_slider_pos(self):
        return self._slider_pos

    def set_slider_pos(self, pos):
        self._slider_pos = pos
        self.update()

    slider_pos = Property(float, get_slider_pos, set_slider_pos)

    def set_current_index(self, index: int):
        if index == self._current_index:
            return
        self._current_index = index
        self._slider_anim.stop()
        self._slider_anim.setStartValue(self._slider_pos)
        self._slider_anim.setEndValue(float(index))
        self._slider_anim.start()
        self.mode_changed.emit(index)

    def current_index(self) -> int:
        return self._current_index

    def minimumSizeHint(self):
        fm = QFontMetrics(QFont("Microsoft YaHei", FONT.MICRO_PT - 1))
        w = 0
        for s in self._labels:
            w += fm.horizontalAdvance(s) + 28
        return QSize(w + 6, 24)

    def sizeHint(self):
        return self.minimumSizeHint()

    def mouseMoveEvent(self, event):
        seg_w = self.width() / len(self._labels)
        idx = int(event.position().x() / seg_w)
        idx = max(0, min(idx, len(self._labels) - 1))
        if idx != self._hovered_segment:
            self._hovered_segment = idx
            self.update()
        super().mouseMoveEvent(event)

    def leaveEvent(self, event):
        self._hovered_segment = -1
        self.update()
        super().leaveEvent(event)

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.setRenderHint(QPainter.RenderHint.Antialiasing)

        w = self.width()
        h = self.height()
        r = BTN.SMALL_BORDER_RADIUS

        painter.setPen(Qt.PenStyle.NoPen)
        painter.setBrush(QColor(COLORS.BG_SECONDARY))
        painter.drawRoundedRect(QRectF(0, 0, w, h), r, r)

        painter.setPen(QPen(QColor(COLORS.BORDER_DEFAULT), 1))
        painter.setBrush(Qt.BrushStyle.NoBrush)
        painter.drawRoundedRect(QRectF(0, 0, w, h), r, r)

        seg_w = w / len(self._labels)
        slider_x = self._slider_pos * seg_w + 2
        slider_w = seg_w - 4
        slider_h = h - 4
        slider_y = 2
        slider_r = max(r - 2, 2)

        painter.setPen(Qt.PenStyle.NoPen)
        painter.setBrush(QColor(COLORS.BRAND))
        painter.drawRoundedRect(QRectF(slider_x, slider_y, slider_w, slider_h), slider_r, slider_r)

        font = QFont()
        font.setPointSize(FONT.MICRO_PT - 1)
        font.setWeight(QFont.Weight.Medium)
        painter.setFont(font)

        for i, seg in enumerate(self._labels):
            if i == self._current_index:
                painter.setPen(QColor(COLORS.BG_PRIMARY))
                font.setBold(True)
                painter.setFont(font)
            elif i == self._hovered_segment:
                painter.setPen(QColor(COLORS.TEXT_PRIMARY))
                font.setBold(False)
                font.setWeight(QFont.Weight.Medium)
                painter.setFont(font)
            else:
                painter.setPen(QColor(COLORS.TEXT_TERTIARY))
                font.setBold(False)
                font.setWeight(QFont.Weight.Medium)
                painter.setFont(font)
            painter.drawText(QRectF(seg_w * i, 0, seg_w, h), Qt.AlignmentFlag.AlignCenter, seg)

        painter.end()

    def mousePressEvent(self, event):
        if event.button() == Qt.MouseButton.LeftButton:
            seg_w = self.width() / len(self._labels)
            idx = int(event.position().x() / seg_w)
            idx = max(0, min(idx, len(self._labels) - 1))
            self.set_current_index(idx)
        super().mousePressEvent(event)


class _ZoomableImageWidget(QWidget):
    """支持缩放和平移的图片预览控件"""
    zoom_changed = Signal(float)

    MIN_ZOOM = 0.1
    MAX_ZOOM = 10.0
    ZOOM_STEP = 0.15

    def __init__(self, parent=None):
        super().__init__(parent)
        self._pixmap = QPixmap()
        self._original_pixmap = QPixmap()
        self._zoom = 1.0
        self._fit_mode = True
        self._fit_zoom = 1.0
        self._offset_x = 0.0
        self._offset_y = 0.0
        self._dragging = False
        self._last_pos = None
        self._border_radius = 6
        self.setSizePolicy(QSizePolicy.Policy.Expanding, QSizePolicy.Policy.Expanding)
        self.setMouseTracking(True)

    def set_pixmap(self, pixmap: QPixmap) -> None:
        """设置要显示的图片"""
        self._pixmap = pixmap
        self._offset_x = 0.0
        self._offset_y = 0.0
        if self._fit_mode:
            self._calc_fit_zoom()
            self._zoom = self._fit_zoom
        self.update()

    def set_fit_mode(self) -> None:
        """切换为适应窗口模式"""
        self._fit_mode = True
        self._calc_fit_zoom()
        self._zoom = self._fit_zoom
        self._offset_x = 0.0
        self._offset_y = 0.0
        self.zoom_changed.emit(self._zoom)
        self.update()

    def set_zoom(self, level: float) -> None:
        """设置缩放级别（0.1 ~ 10.0）"""
        level = max(self.MIN_ZOOM, min(self.MAX_ZOOM, level))
        if abs(level - self._zoom) < 0.001:
            return
        self._fit_mode = False
        self._zoom = level
        self._clamp_offset()
        self.zoom_changed.emit(self._zoom)
        self.update()

    def zoom_level(self) -> float:
        """返回当前缩放级别"""
        return self._zoom

    def zoom_percent(self) -> int:
        """返回当前缩放百分比"""
        return round(self._zoom * 100)

    def is_fit_mode(self) -> bool:
        """是否处于适应窗口模式"""
        return self._fit_mode

    def set_border_radius(self, radius: int) -> None:
        """设置边框圆角"""
        self._border_radius = radius

    def _calc_fit_zoom(self) -> None:
        """计算适应窗口的缩放比例"""
        if self._pixmap.isNull() or self.width() < 10 or self.height() < 10:
            self._fit_zoom = 1.0
            return
        pw = self._pixmap.width()
        ph = self._pixmap.height()
        vw = self.width() - 4
        vh = self.height() - 4
        if pw <= vw and ph <= vh:
            self._fit_zoom = 1.0
        else:
            self._fit_zoom = min(vw / pw, vh / ph)

    def _scaled_size(self) -> QSize:
        """计算缩放后的图片尺寸"""
        if self._pixmap.isNull():
            return QSize(0, 0)
        w = self._pixmap.width() * self._zoom
        h = self._pixmap.height() * self._zoom
        return QSize(int(w), int(h))

    def _clamp_offset(self) -> None:
        """限制偏移量，防止图片拖出可视区域"""
        if self._pixmap.isNull():
            self._offset_x = 0.0
            self._offset_y = 0.0
            return
        ss = self._scaled_size()
        vw = self.width()
        vh = self.height()
        if ss.width() <= vw:
            self._offset_x = 0.0
        else:
            max_x = (ss.width() - vw) / 2.0
            self._offset_x = max(-max_x, min(max_x, self._offset_x))
        if ss.height() <= vh:
            self._offset_y = 0.0
        else:
            max_y = (ss.height() - vh) / 2.0
            self._offset_y = max(-max_y, min(max_y, self._offset_y))

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.setRenderHint(QPainter.RenderHint.Antialiasing)
        painter.setRenderHint(QPainter.RenderHint.SmoothPixmapTransform)

        w = self.width()
        h = self.height()
        r = self._border_radius

        # 绘制背景
        painter.setPen(Qt.PenStyle.NoPen)
        painter.setBrush(QColor(COLORS.BG_PRIMARY))
        painter.drawRoundedRect(QRectF(0, 0, w, h), r, r)

        # 绘制边框
        painter.setPen(QPen(QColor(COLORS.BORDER_DEFAULT), 1))
        painter.setBrush(Qt.BrushStyle.NoBrush)
        painter.drawRoundedRect(QRectF(0, 0, w, h), r, r)

        if self._pixmap.isNull():
            painter.end()
            return

        # 裁剪到圆角区域
        painter.setClipRect(QRectF(1, 1, w - 2, h - 2))

        ss = self._scaled_size()
        # 居中位置 + 偏移
        x = (w - ss.width()) / 2.0 + self._offset_x
        y = (h - ss.height()) / 2.0 + self._offset_y

        target_rect = QRectF(x, y, ss.width(), ss.height())
        source_rect = QRectF(0, 0, self._pixmap.width(), self._pixmap.height())
        painter.drawPixmap(target_rect, self._pixmap, source_rect)
        painter.end()

    def wheelEvent(self, event):
        if event.modifiers() & Qt.KeyboardModifier.ControlModifier:
            delta = event.angleDelta().y()
            if delta > 0:
                factor = 1.0 + self.ZOOM_STEP
            else:
                factor = 1.0 / (1.0 + self.ZOOM_STEP)
            new_zoom = self._zoom * factor
            self.set_zoom(new_zoom)
            event.accept()
        else:
            super().wheelEvent(event)

    def mousePressEvent(self, event):
        if event.button() == Qt.MouseButton.LeftButton:
            ss = self._scaled_size()
            if ss.width() > self.width() or ss.height() > self.height():
                self._dragging = True
                self._last_pos = event.position()
                self.setCursor(Qt.CursorShape.ClosedHandCursor)
                event.accept()
            else:
                super().mousePressEvent(event)
        else:
            super().mousePressEvent(event)

    def mouseMoveEvent(self, event):
        if self._dragging and self._last_pos is not None:
            pos = event.position()
            dx = pos.x() - self._last_pos.x()
            dy = pos.y() - self._last_pos.y()
            self._last_pos = pos
            self._offset_x += dx
            self._offset_y += dy
            self._clamp_offset()
            self.update()
            event.accept()
        else:
            # 更新光标样式
            ss = self._scaled_size()
            if ss.width() > self.width() or ss.height() > self.height():
                self.setCursor(Qt.CursorShape.OpenHandCursor)
            else:
                self.setCursor(Qt.CursorShape.ArrowCursor)
            super().mouseMoveEvent(event)

    def mouseReleaseEvent(self, event):
        if event.button() == Qt.MouseButton.LeftButton and self._dragging:
            self._dragging = False
            self._last_pos = None
            ss = self._scaled_size()
            if ss.width() > self.width() or ss.height() > self.height():
                self.setCursor(Qt.CursorShape.OpenHandCursor)
            else:
                self.setCursor(Qt.CursorShape.ArrowCursor)
            event.accept()
        else:
            super().mouseReleaseEvent(event)

    def resizeEvent(self, event):
        super().resizeEvent(event)
        if self._fit_mode:
            self._calc_fit_zoom()
            self._zoom = self._fit_zoom
            self.zoom_changed.emit(self._zoom)
        else:
            self._clamp_offset()
        self.update()


class PreviewPanel(QWidget):
    preview_requested = Signal()

    def __init__(self, parent=None):
        super().__init__(parent)
        self._current_result = None
        self._full_content = None
        self._showing_truncated = False
        self._current_file_item = None
        self._highlighter = None
        self._md_source_mode = False
        self._html_source_mode = False
        self._md_load_images = True
        self._pdf_doc = None
        self._pdf_total_pages = 0
        self._pdf_pages_loaded = 0
        self._pdf_render_worker = None
        self._pdf_zoom = 1.0
        self._pdf_fit_mode = True
        self._pdf_original_pixmaps: dict[int, QPixmap] = {}  # 存储原始高DPI渲染结果
        self._excel_render_worker = None
        self._preview_timer = QTimer(self)
        self._preview_timer.setSingleShot(True)
        self._preview_timer.setInterval(150)
        self._preview_timer.timeout.connect(self._do_delayed_preview)
        self._pending_file_item = None
        self._pending_force = False
        self._preview_active = False
        self._movie = None
        self._excel_render_mode = False
        self._init_ui()

    def _init_ui(self):
        layout = QVBoxLayout()
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(0)

        _inner_radius = RADIUS.LARGE - 1
        _gap = 4
        _cir = max(RADIUS.LARGE - _gap, 2)

        self.header_widget = QWidget()
        self.header_widget.setObjectName("previewHeader")
        header_layout = QHBoxLayout(self.header_widget)
        header_layout.setContentsMargins(16, 10, 16, 10)
        header_layout.setSpacing(8)

        self.icon_label = QLabel()
        self.icon_label.setFixedSize(20, 20)
        self.icon_label.setStyleSheet("border: none; background: transparent;")

        self.title_label = ElidedFileNameLabel("预览")
        title_font = QFont()
        title_font.setPointSize(FONT.BODY_PT)
        title_font.setBold(True)
        self.title_label.setFont(title_font)
        self.title_label.setStyleSheet(f"color: {COLORS.TEXT_PRIMARY}; border: none; background: transparent;")

        self._file_info_label = QLabel("")
        self._file_info_label.setStyleSheet(f"""
            color: {COLORS.TEXT_TERTIARY};
            font-size: {FONT.MICRO_PT}px;
            border: none;
            background: transparent;
        """)

        header_layout.addWidget(self.icon_label)
        header_layout.addWidget(self.title_label, 1)
        header_layout.addWidget(self._file_info_label)

        self.header_widget.setStyleSheet(f"""
            #previewHeader {{
                background-color: {COLORS.BG_SECONDARY};
                border-bottom: 1px solid {COLORS.BORDER_DEFAULT};
                border-top-left-radius: {_inner_radius}px;
                border-top-right-radius: {_inner_radius}px;
            }}
        """)

        self.empty_placeholder = QWidget()
        self.empty_placeholder.setSizePolicy(QSizePolicy.Policy.Expanding, QSizePolicy.Policy.Expanding)
        empty_layout = QVBoxLayout(self.empty_placeholder)
        empty_layout.setContentsMargins(0, 0, 0, 0)
        empty_layout.setSpacing(8)

        self._empty_icon = QLabel()
        self._empty_icon.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self._empty_icon.setPixmap(QIcon("icons/doctype/File.svg").pixmap(QSize(56, 56)))
        self._empty_icon.setMinimumSize(56, 56)
        self._empty_icon.setStyleSheet("border: none; background: transparent;")

        self._empty_title = QLabel("选择文件以预览")
        self._empty_title.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self._empty_title.setStyleSheet(f"""
            color: {COLORS.TEXT_PLACEHOLDER};
            font-size: {FONT.BODY_PT}px;
            border: none; background: transparent;
        """)

        self._empty_hint = QLabel("点击搜索结果中的文件查看预览")
        self._empty_hint.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self._empty_hint.setStyleSheet(f"""
            color: {COLORS.TEXT_TERTIARY};
            font-size: {FONT.CAPTION_PT}px;
            border: none; background: transparent;
        """)

        # 使用居中容器确保垂直居中
        center_container = QWidget()
        center_container.setSizePolicy(QSizePolicy.Policy.Expanding, QSizePolicy.Policy.Expanding)
        center_layout = QVBoxLayout(center_container)
        center_layout.setContentsMargins(0, 0, 0, 0)
        center_layout.setSpacing(8)
        center_layout.addStretch(1)
        center_layout.addWidget(self._empty_icon, 0, Qt.AlignmentFlag.AlignCenter)
        center_layout.addWidget(self._empty_title, 0, Qt.AlignmentFlag.AlignCenter)
        center_layout.addWidget(self._empty_hint, 0, Qt.AlignmentFlag.AlignCenter)
        center_layout.addStretch(1)

        empty_layout.addWidget(center_container)

        self._content_stack = QStackedWidget()

        self._text_preview = QWidget()
        tp_layout = QVBoxLayout(self._text_preview)
        tp_layout.setContentsMargins(12, 8, 12, 12)
        tp_layout.setSpacing(4)

        self._text_edit = QTextEdit()
        self._text_edit.setReadOnly(True)
        self._text_edit.setStyleSheet(f"""
            QTextEdit {{
                background-color: {COLORS.BG_PRIMARY};
                border: 1px solid {COLORS.BORDER_DEFAULT};
                border-radius: {_cir}px;
                color: {COLORS.TEXT_SECONDARY};
                font-family: "Cascadia Code", "Consolas", "Courier New", monospace;
                font-size: {FONT.MICRO_PT}px;
                selection-background-color: {COLORS.BRAND_LIGHT_BG};
                selection-color: {COLORS.TEXT_PRIMARY};
            }}
        """)
        self._text_edit.setFrameStyle(QTextEdit.Shape.NoFrame)
        self._text_edit.setViewportMargins(16, 12, 16, 12)

        self._show_more_btn = QPushButton("显示剩余内容")
        self._show_more_btn.setCursor(Qt.CursorShape.PointingHandCursor)
        self._show_more_btn.setStyleSheet(button_small_secondary())
        self._show_more_btn.setVisible(False)
        self._show_more_btn.clicked.connect(self._on_show_more)
        # 绘制下拉箭头图标
        _chevron_pixmap = QPixmap(12, 12)
        _chevron_pixmap.fill(Qt.GlobalColor.transparent)
        _chevron_painter = QPainter(_chevron_pixmap)
        _chevron_painter.setRenderHint(QPainter.RenderHint.Antialiasing)
        _chevron_pen = QPen(QColor(COLORS.TEXT_SECONDARY), 1.5)
        _chevron_pen.setCapStyle(Qt.PenCapStyle.RoundCap)
        _chevron_pen.setJoinStyle(Qt.PenJoinStyle.RoundJoin)
        _chevron_painter.setPen(_chevron_pen)
        _chevron_painter.drawLine(2, 4, 6, 8)
        _chevron_painter.drawLine(6, 8, 10, 4)
        _chevron_painter.end()
        self._show_more_btn.setIcon(QIcon(_chevron_pixmap))

        self._md_mode_slider = _PreviewModeSlider(['预览', '源码'])
        self._md_mode_slider.setVisible(False)
        self._md_mode_slider.mode_changed.connect(self._on_toggle_md_mode)

        self._html_mode_slider = _PreviewModeSlider(['渲染', '源码'])
        self._html_mode_slider.setVisible(False)
        self._html_mode_slider.mode_changed.connect(self._on_toggle_html_mode)

        btn_row = QHBoxLayout()
        btn_row.addStretch()
        btn_row.addWidget(self._html_mode_slider)
        btn_row.addWidget(self._md_mode_slider)
        btn_row.addWidget(self._show_more_btn)

        self._md_browser = QTextBrowser()
        self._md_browser.setOpenExternalLinks(False)
        self._md_browser.setStyleSheet(f"""
            QTextBrowser {{
                background-color: {COLORS.BG_PRIMARY};
                border: 1px solid {COLORS.BORDER_DEFAULT};
                border-radius: {_cir}px;
                color: {COLORS.TEXT_SECONDARY};
                font-size: {FONT.CAPTION_PT}px;
                selection-background-color: {COLORS.BRAND_LIGHT_BG};
            }}
        """)
        self._md_browser.setViewportMargins(20, 16, 20, 16)
        self._md_browser.setVisible(False)

        tp_layout.addWidget(self._md_browser, 1)
        tp_layout.addWidget(self._text_edit, 1)
        tp_layout.addLayout(btn_row)

        self._image_preview = QWidget()
        ip_layout = QVBoxLayout(self._image_preview)
        ip_layout.setContentsMargins(12, 8, 12, 12)
        ip_layout.setSpacing(4)

        self._image_label = _ZoomableImageWidget()
        self._image_label.set_border_radius(_cir)
        self._image_label.zoom_changed.connect(self._on_image_zoom_changed)

        # GIF 模式用的 QLabel（隐藏，仅 GIF 时显示）
        self._gif_label = QLabel()
        self._gif_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self._gif_label.setSizePolicy(QSizePolicy.Policy.Expanding, QSizePolicy.Policy.Expanding)
        self._gif_label.setStyleSheet(f"""
            border: 1px solid {COLORS.BORDER_DEFAULT};
            border-radius: {_cir}px;
            background-color: {COLORS.BG_PRIMARY};
        """)
        self._gif_label.setVisible(False)

        self._fit_view_btn = QPushButton("适应窗口")
        self._fit_view_btn.setCursor(Qt.CursorShape.PointingHandCursor)
        self._fit_view_btn.setStyleSheet(button_small_secondary())
        self._fit_view_btn.setVisible(False)
        self._fit_view_btn.clicked.connect(self._on_fit_view)

        self._zoom_slider = QSlider(Qt.Orientation.Horizontal)
        self._zoom_slider.setFixedWidth(120)
        self._zoom_slider.setMinimumHeight(22)
        self._zoom_slider.setMinimum(10)
        self._zoom_slider.setMaximum(1000)
        self._zoom_slider.setValue(100)
        self._zoom_slider.setStyleSheet(f"""
            QSlider::groove:horizontal {{
                border: none;
                height: 4px;
                background: {COLORS.BG_TERTIARY};
                border-radius: 2px;
            }}
            QSlider::handle:horizontal {{
                background: {COLORS.BRAND};
                border: 2px solid {COLORS.BG_PRIMARY};
                width: 14px;
                height: 14px;
                margin: -5px 0;
                border-radius: 7px;
            }}
            QSlider::sub-page:horizontal {{
                background: {COLORS.BRAND};
                border-radius: 2px;
            }}
        """)
        self._zoom_slider.setVisible(False)
        self._zoom_slider.valueChanged.connect(self._on_zoom_slider_changed)

        self._zoom_percent_label = QLabel("100%")
        self._zoom_percent_label.setFixedWidth(42)
        self._zoom_percent_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self._zoom_percent_label.setStyleSheet(f"""
            color: {COLORS.TEXT_TERTIARY};
            font-size: {FONT.MICRO_PT}px;
            border: none;
            background: transparent;
        """)
        self._zoom_percent_label.setVisible(False)

        self._view_original_btn = QPushButton("查看原图")
        self._view_original_btn.setCursor(Qt.CursorShape.PointingHandCursor)
        self._view_original_btn.setStyleSheet(button_small_primary())
        self._view_original_btn.setVisible(False)
        self._view_original_btn.clicked.connect(self._on_view_original_image)

        ip_btn_row = QHBoxLayout()
        ip_btn_row.setContentsMargins(0, 0, 0, 0)
        ip_btn_row.addStretch()
        ip_btn_row.addWidget(self._fit_view_btn)
        ip_btn_row.addWidget(self._zoom_percent_label)
        ip_btn_row.addWidget(self._zoom_slider)
        ip_btn_row.addWidget(self._view_original_btn)

        # 按钮行容器，固定高度避免滑块拖动时布局跳动
        ip_btn_container = QWidget()
        ip_btn_container.setFixedHeight(34)
        ip_btn_container.setLayout(ip_btn_row)

        # 使用堆叠方式：_image_label 和 _gif_label 占据同一位置
        from PySide6.QtWidgets import QStackedLayout
        image_stack = QStackedLayout()
        image_stack.setStackingMode(QStackedLayout.StackingMode.StackAll)
        image_stack.addWidget(self._image_label)
        image_stack.addWidget(self._gif_label)
        # 创建一个容器 widget
        image_container = QWidget()
        image_container.setLayout(image_stack)

        ip_layout.addWidget(image_container, 1)
        ip_layout.addWidget(ip_btn_container)

        self._pdf_preview = QWidget()
        pdf_layout = QVBoxLayout(self._pdf_preview)
        pdf_layout.setContentsMargins(12, 8, 12, 12)
        pdf_layout.setSpacing(4)

        self._pdf_scroll_area = QScrollArea()
        self._pdf_scroll_area.setWidgetResizable(True)
        self._pdf_scroll_area.setStyleSheet(f"""
            QScrollArea {{
                background-color: {COLORS.BG_PRIMARY};
                border: 1px solid {COLORS.BORDER_DEFAULT};
                border-radius: {_cir}px;
            }}
        """ + scrollbar_style())
        self._pdf_scroll_area.setFrameStyle(QScrollArea.Shape.NoFrame)
        self._pdf_scroll_area.wheelEvent = self._on_pdf_scroll_wheel

        # PDF 页面容器
        self._pdf_pages_widget = QWidget()
        self._pdf_pages_layout = QVBoxLayout(self._pdf_pages_widget)
        self._pdf_pages_layout.setContentsMargins(0, 0, 0, 0)
        self._pdf_pages_layout.setSpacing(8)
        self._pdf_pages_layout.setAlignment(Qt.AlignmentFlag.AlignTop | Qt.AlignmentFlag.AlignHCenter)
        self._pdf_page_labels: list[QLabel] = []
        self._pdf_pages_widget.setStyleSheet(f"background-color: {COLORS.BG_PRIMARY};")
        self._pdf_scroll_area.setWidget(self._pdf_pages_widget)

        # PDF 缩放控件
        self._pdf_fit_btn = QPushButton("适应窗口")
        self._pdf_fit_btn.setCursor(Qt.CursorShape.PointingHandCursor)
        self._pdf_fit_btn.setStyleSheet(button_small_secondary())
        self._pdf_fit_btn.setVisible(False)
        self._pdf_fit_btn.clicked.connect(self._on_pdf_fit_view)

        self._pdf_zoom_label = QLabel("100%")
        self._pdf_zoom_label.setFixedWidth(42)
        self._pdf_zoom_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self._pdf_zoom_label.setStyleSheet(f"""
            color: {COLORS.TEXT_TERTIARY};
            font-size: {FONT.MICRO_PT}px;
            border: none;
            background: transparent;
        """)
        self._pdf_zoom_label.setVisible(False)

        self._pdf_zoom_slider = QSlider(Qt.Orientation.Horizontal)
        self._pdf_zoom_slider.setFixedWidth(80)
        self._pdf_zoom_slider.setMinimumHeight(22)
        self._pdf_zoom_slider.setMinimum(25)
        self._pdf_zoom_slider.setMaximum(500)
        self._pdf_zoom_slider.setValue(100)
        self._pdf_zoom_slider.setStyleSheet(f"""
            QSlider::groove:horizontal {{
                border: none;
                height: 4px;
                background: {COLORS.BG_TERTIARY};
                border-radius: 2px;
            }}
            QSlider::handle:horizontal {{
                background: {COLORS.BRAND};
                border: 2px solid {COLORS.BG_PRIMARY};
                width: 14px;
                height: 14px;
                margin: -5px 0;
                border-radius: 7px;
            }}
            QSlider::sub-page:horizontal {{
                background: {COLORS.BRAND};
                border-radius: 2px;
            }}
        """)
        self._pdf_zoom_slider.setVisible(False)
        self._pdf_zoom_slider.valueChanged.connect(self._on_pdf_zoom_slider_changed)

        self._pdf_load_more_btn = QPushButton("加载更多页面")
        self._pdf_load_more_btn.setCursor(Qt.CursorShape.PointingHandCursor)
        self._pdf_load_more_btn.setStyleSheet(button_small_primary())
        self._pdf_load_more_btn.setVisible(False)
        self._pdf_load_more_btn.clicked.connect(self._on_pdf_load_more)

        pdf_btn_row = QHBoxLayout()
        pdf_btn_row.setContentsMargins(0, 0, 0, 0)
        pdf_btn_row.addStretch()
        pdf_btn_row.addWidget(self._pdf_fit_btn)
        pdf_btn_row.addWidget(self._pdf_zoom_label)
        pdf_btn_row.addWidget(self._pdf_zoom_slider)
        pdf_btn_row.addWidget(self._pdf_load_more_btn)

        # 按钮行容器，固定高度避免滑块拖动时布局跳动
        pdf_btn_container = QWidget()
        pdf_btn_container.setFixedHeight(34)
        pdf_btn_container.setLayout(pdf_btn_row)

        pdf_layout.addWidget(self._pdf_scroll_area, 1)
        pdf_layout.addWidget(pdf_btn_container)

        self._excel_preview = QWidget()
        excel_layout = QVBoxLayout(self._excel_preview)
        excel_layout.setContentsMargins(12, 8, 12, 12)
        excel_layout.setSpacing(4)

        self._excel_tab = QTabWidget()
        self._excel_tab.setStyleSheet(f"""
            QTabWidget::pane {{
                border: 1px solid {COLORS.BORDER_DEFAULT};
                border-radius: {_cir}px;
                background-color: {COLORS.BG_PRIMARY};
            }}
            QTabBar::tab {{
                background-color: {COLORS.BG_SECONDARY};
                color: {COLORS.TEXT_TERTIARY};
                border: 1px solid {COLORS.BORDER_DEFAULT};
                border-bottom: none;
                border-top-left-radius: {RADIUS.SMALL}px;
                border-top-right-radius: {RADIUS.SMALL}px;
                padding: 4px 12px;
                margin-right: 2px;
                font-size: {FONT.MICRO_PT}px;
            }}
            QTabBar::tab:selected {{
                background-color: {COLORS.BG_PRIMARY};
                color: {COLORS.TEXT_PRIMARY};
                font-weight: bold;
            }}
            QTabBar::tab:hover {{
                background-color: {COLORS.BG_TERTIARY};
            }}
        """)

        self._excel_mode_slider = _PreviewModeSlider(['文本预览', '表格预览'])
        self._excel_mode_slider.mode_changed.connect(self._on_toggle_excel_mode)

        excel_btn_row = QHBoxLayout()
        excel_btn_row.addStretch()
        excel_btn_row.addWidget(self._excel_mode_slider)

        self._excel_text_edit = QTextEdit()
        self._excel_text_edit.setReadOnly(True)
        self._excel_text_edit.setStyleSheet(f"""
            QTextEdit {{
                background-color: {COLORS.BG_PRIMARY};
                border: 1px solid {COLORS.BORDER_DEFAULT};
                border-radius: {_cir}px;
                color: {COLORS.TEXT_TERTIARY};
                font-family: "Cascadia Code", "Consolas", "Courier New", monospace;
                font-size: {FONT.MICRO_PT - 1}px;
            }}
        """)
        self._excel_text_edit.setFrameStyle(QTextEdit.Shape.NoFrame)
        self._excel_text_edit.setViewportMargins(12, 8, 12, 8)

        excel_layout.addWidget(self._excel_tab, 1)
        excel_layout.addWidget(self._excel_text_edit, 1)
        excel_layout.addLayout(excel_btn_row)

        self._media_preview = QWidget()
        mp_layout = QVBoxLayout(self._media_preview)
        mp_layout.setContentsMargins(12, 8, 12, 12)
        mp_layout.setSpacing(4)

        # 媒体属性信息预览区域（类似右键属性对话框）
        self._media_fallback = QWidget()
        self._media_fallback.setStyleSheet(f"""
            QWidget {{
                background-color: {COLORS.BG_SECONDARY};
                border: 1px solid {COLORS.BORDER_DEFAULT};
                border-radius: {_cir}px;
            }}
        """)
        fallback_layout = QVBoxLayout(self._media_fallback)
        fallback_layout.setAlignment(Qt.AlignmentFlag.AlignTop)
        fallback_layout.setContentsMargins(24, 20, 24, 20)
        fallback_layout.setSpacing(12)

        # 封面图片/缩略图
        self._media_cover_label = QLabel()
        self._media_cover_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self._media_cover_label.setStyleSheet(f"""
            border: 1px solid {COLORS.BORDER_DEFAULT};
            border-radius: {RADIUS.MEDIUM}px;
            background-color: #1a1a2e;
        """)
        self._media_cover_label.setFixedSize(280, 200)
        self._media_cover_label.setScaledContents(False)
        self._media_cover_label.setVisible(False)

        # 详细信息表格
        self._media_info_table = QTableWidget()
        self._media_info_table.setEditTriggers(QTableWidget.EditTrigger.NoEditTriggers)
        self._media_info_table.setSelectionMode(QTableWidget.SelectionMode.NoSelection)
        self._media_info_table.verticalHeader().setVisible(False)
        self._media_info_table.horizontalHeader().setVisible(False)
        self._media_info_table.horizontalHeader().setStretchLastSection(True)
        self._media_info_table.setShowGrid(False)
        self._media_info_table.setStyleSheet(f"""
            QTableWidget {{
                background-color: transparent;
                border: none;
                font-size: {FONT.MICRO_PT - 1}px;
            }}
            QTableWidget::item {{
                padding: 2px 6px;
                border-bottom: 1px solid {COLORS.BORDER_DEFAULT};
                border-right: none;
                border-left: none;
                border-top: none;
            }}
        """)
        self._media_info_table.setVisible(False)

        fallback_layout.addWidget(self._media_cover_label, 0, Qt.AlignmentFlag.AlignCenter)
        fallback_layout.addWidget(self._media_info_table)

        self._media_fallback.setVisible(False)

        self._open_system_player_btn = QPushButton("使用系统播放器打开")
        self._open_system_player_btn.setCursor(Qt.CursorShape.PointingHandCursor)
        self._open_system_player_btn.setStyleSheet(button_small_secondary())
        self._open_system_player_btn.setVisible(False)
        self._open_system_player_btn.clicked.connect(self._on_open_with_system_player)

        mp_btn_row = QHBoxLayout()
        mp_btn_row.addStretch()
        mp_btn_row.addWidget(self._open_system_player_btn)

        mp_layout.addWidget(self._media_fallback, 1)
        mp_layout.addLayout(mp_btn_row)

        self._unsupported_preview = QWidget()
        us_layout = QVBoxLayout(self._unsupported_preview)
        us_layout.setContentsMargins(12, 8, 12, 12)
        us_layout.setSpacing(0)

        us_box = QWidget()
        us_box.setStyleSheet(f"""
            QWidget {{
                background-color: {COLORS.BG_SECONDARY};
                border: 1px solid {COLORS.BORDER_DEFAULT};
                border-radius: {_cir}px;
            }}
        """)
        us_box_layout = QVBoxLayout(us_box)
        us_box_layout.setContentsMargins(24, 24, 24, 24)
        us_box_layout.setAlignment(Qt.AlignmentFlag.AlignCenter)
        us_box_layout.setSpacing(8)

        self._unsupported_icon = QLabel()
        self._unsupported_icon.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self._unsupported_icon.setStyleSheet("border: none; background: transparent;")

        self._unsupported_text = QLabel("此文件类型暂不支持预览")
        self._unsupported_text.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self._unsupported_text.setStyleSheet(f"""
            color: {COLORS.TEXT_PLACEHOLDER};
            font-size: {FONT.BODY_PT}px;
            border: none; background: transparent;
        """)

        self._unsupported_hint = QLabel("")
        self._unsupported_hint.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self._unsupported_hint.setWordWrap(True)
        self._unsupported_hint.setStyleSheet(f"""
            color: {COLORS.TEXT_TERTIARY};
            font-size: {FONT.CAPTION_PT}px;
            border: none; background: transparent;
        """)

        self._force_preview_btn = QPushButton("仍然预览")
        self._force_preview_btn.setCursor(Qt.CursorShape.PointingHandCursor)
        self._force_preview_btn.setStyleSheet(button_small_primary())
        self._force_preview_btn.setVisible(False)
        self._force_preview_btn.clicked.connect(self._on_force_preview)

        us_box_layout.addWidget(self._unsupported_icon)
        us_box_layout.addWidget(self._unsupported_text)
        us_box_layout.addWidget(self._unsupported_hint)
        us_box_layout.addWidget(self._force_preview_btn)

        us_layout.addStretch()
        us_layout.addWidget(us_box, 1)
        us_layout.addStretch()

        self._folder_preview = QWidget()
        fp_layout = QVBoxLayout(self._folder_preview)
        fp_layout.setContentsMargins(12, 8, 12, 12)
        fp_layout.setSpacing(6)

        self._folder_stats_label = QLabel("")
        self._folder_stats_label.setStyleSheet(label_caption_style())

        self._folder_scroll = QScrollArea()
        self._folder_scroll.setWidgetResizable(True)
        self._folder_scroll.setHorizontalScrollBarPolicy(Qt.ScrollBarPolicy.ScrollBarAlwaysOff)
        self._folder_scroll.setVerticalScrollBarPolicy(Qt.ScrollBarPolicy.ScrollBarAsNeeded)
        self._folder_scroll.setStyleSheet(
            f"QScrollArea {{ border: 1px solid {COLORS.BORDER_DEFAULT}; border-radius: {_cir}px; background-color: {COLORS.BG_SECONDARY}; }}" + scrollbar_style()
        )

        self._folder_list_container = QWidget()
        self._folder_list_layout = QVBoxLayout(self._folder_list_container)
        self._folder_list_layout.setContentsMargins(4, 4, 4, 4)
        self._folder_list_layout.setSpacing(0)
        self._folder_list_layout.setAlignment(Qt.AlignmentFlag.AlignTop)
        self._folder_list_container.setStyleSheet("QWidget { background: transparent; }")
        self._folder_scroll.setWidget(self._folder_list_container)

        fp_layout.addWidget(self._folder_stats_label)
        fp_layout.addWidget(self._folder_scroll, 1)

        self._loading_preview = QWidget()
        ll_layout = QVBoxLayout(self._loading_preview)
        ll_layout.setContentsMargins(12, 8, 12, 12)
        ll_layout.setSpacing(0)
        ll_layout.setAlignment(Qt.AlignmentFlag.AlignCenter)

        ll_box = QWidget()
        ll_box.setStyleSheet(f"""
            QWidget {{
                background-color: {COLORS.BG_SECONDARY};
                border: 1px solid {COLORS.BORDER_DEFAULT};
                border-radius: {_cir}px;
            }}
        """)
        ll_box_layout = QVBoxLayout(ll_box)
        ll_box_layout.setContentsMargins(24, 24, 24, 24)
        ll_box_layout.setAlignment(Qt.AlignmentFlag.AlignCenter)
        ll_box_layout.setSpacing(12)

        self._loading_spinner = _PreviewLoadingSpinner(ll_box)
        self._loading_text = QLabel("正在加载预览...")
        self._loading_text.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self._loading_text.setStyleSheet(f"""
            color: {COLORS.TEXT_TERTIARY};
            font-size: {FONT.CAPTION_PT}px;
            border: none; background: transparent;
        """)

        self._loading_progress = QProgressBar()
        self._loading_progress.setFixedHeight(4)
        self._loading_progress.setMinimum(0)
        self._loading_progress.setMaximum(0)
        self._loading_progress.setTextVisible(False)
        self._loading_progress.setStyleSheet(f"""
            QProgressBar {{
                background-color: {COLORS.BG_TERTIARY};
                border-radius: 2px;
                border: none;
            }}
            QProgressBar::chunk {{
                background-color: {COLORS.BRAND};
                border-radius: 2px;
            }}
        """)
        self._loading_progress.setVisible(False)

        ll_box_layout.addWidget(self._loading_spinner, 0, Qt.AlignmentFlag.AlignCenter)
        ll_box_layout.addWidget(self._loading_text)
        ll_box_layout.addWidget(self._loading_progress)

        ll_layout.addStretch()
        ll_layout.addWidget(ll_box, 0)
        ll_layout.addStretch()

        self._ppt_preview = QWidget()
        ppt_layout = QVBoxLayout(self._ppt_preview)
        ppt_layout.setContentsMargins(12, 8, 12, 12)
        ppt_layout.setSpacing(4)

        self._ppt_text_edit = QTextEdit()
        self._ppt_text_edit.setReadOnly(True)
        self._ppt_text_edit.setStyleSheet(f"""
            QTextEdit {{
                background-color: {COLORS.BG_PRIMARY};
                border: 1px solid {COLORS.BORDER_DEFAULT};
                border-radius: {_cir}px;
                color: {COLORS.TEXT_SECONDARY};
                font-family: "Microsoft YaHei", sans-serif;
                font-size: {FONT.CAPTION_PT}px;
            }}
        """)
        self._ppt_text_edit.setFrameStyle(QTextEdit.Shape.NoFrame)
        self._ppt_text_edit.setViewportMargins(16, 12, 16, 12)

        ppt_layout.addWidget(self._ppt_text_edit, 1)

        self._content_stack.addWidget(self._text_preview)
        self._content_stack.addWidget(self._image_preview)
        self._content_stack.addWidget(self._pdf_preview)
        self._content_stack.addWidget(self._excel_preview)
        self._content_stack.addWidget(self._ppt_preview)
        self._content_stack.addWidget(self._media_preview)
        self._content_stack.addWidget(self._unsupported_preview)
        self._content_stack.addWidget(self._folder_preview)
        self._content_stack.addWidget(self._loading_preview)
        self._content_stack.addWidget(self.empty_placeholder)

        layout.addWidget(self.header_widget)
        layout.addWidget(self._content_stack, 1)

        self.setLayout(layout)
        self.setStyleSheet(f"""
            PreviewPanel {{
                background-color: {COLORS.BG_PRIMARY};
            }}
        """)

    def _should_auto_preview(self, file_item) -> bool:
        if file_item.is_directory:
            return True
        ext = file_item.extension.lower()
        if ext in ARCHIVE_EXTENSIONS:
            return True
        return False

    def set_result(self, result):
        self._cancel_workers()
        self._stop_media()
        self._stop_gif()
        # 释放图片资源
        if hasattr(self._image_label, '_original_pixmap'):
            self._image_label._original_pixmap = QPixmap()
        self._image_label.set_pixmap(QPixmap())
        self._text_edit.clear()
        self._md_browser.clear()
        self._excel_text_edit.clear()
        self._current_result = result
        self._full_content = None
        self._showing_truncated = False
        self._current_file_item = None
        self._highlighter = None
        self._md_source_mode = False
        self._pdf_doc = None
        self._pdf_total_pages = 0
        self._pdf_pages_loaded = 0
        self._preview_active = False
        self._excel_render_mode = False
        if result is None:
            self._show_empty()
            return
        file_item = result.file_item
        if file_item.is_directory:
            self._show_folder_preview(file_item)
        elif self._should_auto_preview(file_item):
            self._show_file_preview(file_item)
        else:
            self._show_preview_prompt(file_item)

    def _show_preview_prompt(self, file_item):
        self._hide_all_content()
        self._current_file_item = file_item
        self.title_label.setText(file_item.name)
        self._file_info_label.setText(file_item.size_display)

        ext = file_item.extension.lower()
        icon_name = FILE_ICON_MAP.get(ext, 'doctype/File.svg')
        self.icon_label.setPixmap(QIcon(f"icons/{icon_name}").pixmap(QSize(20, 20)))

        self._unsupported_icon.setPixmap(
            QIcon(f"icons/{icon_name}").pixmap(QSize(48, 48))
        )
        self._unsupported_text.setText("按空格键预览")
        self._unsupported_hint.setText(
            f"点击此处或按空格键开启预览\n可双击文件使用默认程序打开"
        )
        self._force_preview_btn.setVisible(True)
        self._force_preview_btn.setText("预览文件")
        self._content_stack.setCurrentWidget(self._unsupported_preview)

    def activate_preview(self):
        if self._current_file_item and not self._preview_active:
            self._show_file_preview(self._current_file_item)

    def _show_empty(self):
        self._hide_all_content()
        self.title_label.setText("预览")
        self.icon_label.clear()
        self._file_info_label.setText("")
        self._content_stack.setCurrentWidget(self.empty_placeholder)

    def _hide_all_content(self):
        self._show_more_btn.setVisible(False)
        self._md_mode_slider.setVisible(False)
        self._html_mode_slider.setVisible(False)
        self._view_original_btn.setVisible(False)
        self._force_preview_btn.setVisible(False)
        self._pdf_load_more_btn.setVisible(False)
        self._pdf_fit_btn.setVisible(False)
        self._pdf_zoom_label.setVisible(False)
        self._pdf_zoom_slider.setVisible(False)
        self._loading_spinner.stop()
        self._loading_progress.setVisible(False)
        self._zoom_slider.setVisible(False)
        self._zoom_percent_label.setVisible(False)
        self._fit_view_btn.setVisible(False)

    def _get_dpi_for_widget(self) -> int:
        screen = self.screen()
        if screen:
            return int(screen.logicalDotsPerInch() * 1.5)
        return 150

    def _show_file_preview(self, file_item, force=False):
        self._hide_all_content()
        self._current_file_item = file_item
        self._preview_active = True
        self.title_label.setText(file_item.name)
        self._file_info_label.setText(file_item.size_display)

        ext = file_item.extension.lower()
        icon_name = FILE_ICON_MAP.get(ext, 'doctype/File.svg')
        self.icon_label.setPixmap(QIcon(f"icons/{icon_name}").pixmap(QSize(20, 20)))

        if ext in MARKDOWN_EXTS:
            self._show_markdown_content(file_item, force)
        elif ext in HTML_EXTS:
            self._show_html_content(file_item, force)
        elif ext in PREVIEWABLE_TEXT_EXTS:
            self._show_text_content(file_item, force)
        elif ext in PREVIEWABLE_IMAGE_EXTS:
            self._show_image_content(file_item)
        elif ext == '.pdf':
            self._show_pdf_content(file_item, force)
        elif ext in {'.docx', '.doc'}:
            self._show_word_content(file_item, force)
        elif ext in {'.xlsx', '.xls'}:
            self._show_excel_content(file_item, force)
        elif ext in {'.pptx', '.ppt'}:
            self._show_ppt_content(file_item, force)
        elif ext in VIDEO_EXTENSIONS:
            self._show_video_content(file_item, force)
        elif ext in AUDIO_EXTENSIONS:
            self._show_audio_content(file_item, force)
        elif ext in ARCHIVE_EXTENSIONS:
            self._show_archive_content(file_item)
        else:
            self._show_unsupported(file_item, "此文件类型暂不支持预览\n可双击文件使用默认程序打开查看")

    def _check_file_size(self, file_item, max_mb, force, hint_suffix=""):
        try:
            file_size = os.path.getsize(file_item.path)
        except OSError:
            self._show_unsupported(file_item, "无法读取文件")
            return None
        if file_size > max_mb * 1024 * 1024 and not force:
            self._show_oversized(file_item, max_mb, hint_suffix)
            return None
        return file_size

    def _show_oversized(self, file_item, max_mb, hint_suffix=""):
        self._unsupported_icon.setPixmap(
            QIcon(f"icons/{FILE_ICON_MAP.get(file_item.extension.lower(), 'doctype/File.svg')}").pixmap(QSize(48, 48))
        )
        self._unsupported_text.setText("文件过大")
        self._unsupported_hint.setText(
            f"文件大小 {file_item.size_display}，超过预览限制 {max_mb}MB{hint_suffix}\n可点击下方按钮尝试预览"
        )
        self._force_preview_btn.setVisible(True)
        self._force_preview_btn.setText("仍然预览")
        self._content_stack.setCurrentWidget(self._unsupported_preview)

    def _show_text_content(self, file_item, force=False):
        file_size = self._check_file_size(file_item, MAX_TEXT_PREVIEW_SIZE_MB, force)
        if file_size is None and not force:
            return

        from utils.encoding import read_text_file
        max_mb = MAX_TEXT_PREVIEW_SIZE_MB if not force else 50
        content = read_text_file(file_item.path, max_size_mb=max_mb)
        if content is None:
            self._show_unsupported(file_item, "无法读取文件内容")
            return

        self._full_content = content
        lines = content.splitlines()
        truncated = len(lines) > MAX_PREVIEW_LINES
        display_lines = lines[:MAX_PREVIEW_LINES] if truncated else lines
        self._showing_truncated = truncated

        display_content = '\n'.join(display_lines)
        self._text_edit.setPlainText(display_content)

        ext = file_item.extension.lower()
        self._highlighter = _create_highlighter(self._text_edit.document(), ext)

        cursor = self._text_edit.textCursor()
        cursor.movePosition(QTextCursor.MoveOperation.Start)
        self._text_edit.setTextCursor(cursor)

        self._show_more_btn.setVisible(truncated)
        self._md_mode_slider.setVisible(False)
        self._html_mode_slider.setVisible(False)
        self._md_browser.setVisible(False)
        self._text_edit.setVisible(True)

        if truncated:
            self._file_info_label.setText(
                f"{file_item.size_display} · 显示前 {MAX_PREVIEW_LINES} 行（共 {len(lines)} 行）"
            )
        else:
            self._file_info_label.setText(f"{file_item.size_display} · {len(lines)} 行")

        self._content_stack.setCurrentWidget(self._text_preview)

    def _show_markdown_content(self, file_item, force=False):
        file_size = self._check_file_size(file_item, MAX_TEXT_PREVIEW_SIZE_MB, force)
        if file_size is None and not force:
            return

        from utils.encoding import read_text_file
        max_mb = MAX_TEXT_PREVIEW_SIZE_MB if not force else 50
        content = read_text_file(file_item.path, max_size_mb=max_mb)
        if content is None:
            self._show_unsupported(file_item, "无法读取文件内容")
            return

        self._full_content = content
        lines = content.splitlines()
        truncated = len(lines) > MAX_PREVIEW_LINES
        self._showing_truncated = truncated

        self._md_source_mode = False
        self._md_mode_slider.setVisible(True)
        self._md_mode_slider.set_current_index(0)
        self._html_mode_slider.setVisible(False)
        self._show_more_btn.setVisible(truncated)

        self._render_markdown(content, truncated)

        if truncated:
            self._file_info_label.setText(
                f"{file_item.size_display} · 预览模式 · 显示前 {MAX_PREVIEW_LINES} 行"
            )
        else:
            self._file_info_label.setText(f"{file_item.size_display} · 预览模式")

        self._content_stack.setCurrentWidget(self._text_preview)

    def _render_markdown(self, content, truncated):
        lines = content.splitlines()
        display_lines = lines[:MAX_PREVIEW_LINES] if truncated else lines
        display_content = '\n'.join(display_lines)

        if self._md_source_mode:
            self._text_edit.setPlainText(display_content)
            self._text_edit.setVisible(True)
            self._md_browser.setVisible(False)
            self._highlighter = _create_highlighter(self._text_edit.document(), '.md')
        else:
            import markdown
            html = markdown.markdown(display_content, extensions=['tables', 'fenced_code', 'codehilite'])
            if not self._md_load_images:
                # 不加载图片时，将 <img> 标签替换为占位文本
                def _replace_img_tag(match):
                    alt_match = re.search(r'alt="([^"]*)"', match.group(0))
                    alt_text = alt_match.group(1) if alt_match else '图片'
                    return f'<span style="color: {COLORS.TEXT_PLACEHOLDER}; font-size: {FONT.MICRO_PT}px; background: {COLORS.BG_TERTIARY}; padding: 2px 6px; border-radius: 3px;">[{alt_text}]</span>'
                html = re.sub(r'<img\s[^>]*/?>', _replace_img_tag, html)
            else:
                # 将相对路径的图片转为绝对 file:// URL
                if self._current_file_item:
                    file_dir = os.path.dirname(os.path.abspath(self._current_file_item.path))
                    def _replace_img_src(match):
                        src = match.group(1)
                        if src and not src.startswith(('http://', 'https://', 'file://', 'data:')):
                            abs_path = os.path.normpath(os.path.join(file_dir, src))
                            file_url = 'file:///' + abs_path.replace('\\', '/')
                            return f'src="{file_url}"'
                        return match.group(0)
                    html = re.sub(r'src="([^"]*)"', _replace_img_src, html)
            styled_html = f"""
            <style>
                body {{ font-family: "Microsoft YaHei", sans-serif; color: {COLORS.TEXT_SECONDARY};
                       font-size: {FONT.CAPTION_PT}px; line-height: 1.6; }}
                h1, h2, h3, h4, h5, h6 {{ color: {COLORS.TEXT_PRIMARY}; margin-top: 12px; margin-bottom: 6px; }}
                h1 {{ font-size: {FONT.DISPLAY_PT}px; }} h2 {{ font-size: {FONT.TITLE_PT}px; }}
                code {{ background-color: {COLORS.BG_TERTIARY}; padding: 2px 6px; border-radius: 3px;
                        font-family: "Cascadia Code", "Consolas", monospace; font-size: {FONT.MICRO_PT}px; }}
                pre {{ background-color: {COLORS.BG_TERTIARY}; padding: 12px; border-radius: 6px;
                       overflow-x: auto; }}
                pre code {{ background: transparent; padding: 0; }}
                blockquote {{ border-left: 3px solid {COLORS.BRAND}; padding-left: 12px;
                             color: {COLORS.TEXT_TERTIARY}; margin: 8px 0; }}
                a {{ color: {COLORS.BRAND}; text-decoration: none; }}
                table {{ border-collapse: collapse; margin: 8px 0; }}
                th, td {{ border: 1px solid {COLORS.BORDER_DEFAULT}; padding: 6px 12px; }}
                th {{ background-color: {COLORS.BG_TERTIARY}; }}
                ul, ol {{ padding-left: 20px; }}
                img {{ max-width: 100%; }}
            </style>
            {html}
            """
            self._md_browser.setHtml(styled_html)
            self._md_browser.setVisible(True)
            self._text_edit.setVisible(False)

    def _show_image_content(self, file_item):
        try:
            file_size = os.path.getsize(file_item.path)
        except OSError:
            self._show_unsupported(file_item, "无法读取文件")
            return

        if file_size > 50 * 1024 * 1024:
            self._show_oversized(file_item, 50)
            return

        ext = file_item.extension.lower()
        if ext == '.gif':
            self._show_gif_content(file_item)
            return

        # 非 GIF 模式：显示 _image_label，隐藏 _gif_label
        self._image_label.setVisible(True)
        self._gif_label.setVisible(False)

        pixmap = QPixmap()
        try:
            from PySide6.QtGui import QImageReader
            reader = QImageReader(file_item.path)
            reader.setAllocationLimit(0)
            img = reader.read()
            if img is not None and not img.isNull():
                pixmap = QPixmap.fromImage(img)
            else:
                pixmap = QPixmap(file_item.path)
        except Exception:
            pixmap = QPixmap(file_item.path)

        if pixmap.isNull():
            self._show_unsupported(file_item, "无法加载图片\n可能是图片过大或格式不受支持")
            return

        # 如果图片像素过多，先缩放到合理尺寸再用于预览
        original_pixmap = pixmap
        is_scaled = False
        if pixmap.width() * pixmap.height() > MAX_IMAGE_PREVIEW_PIXELS:
            self._image_label._original_pixmap = pixmap
            scaled = pixmap.scaled(
                1920, 1080,
                Qt.AspectRatioMode.KeepAspectRatio,
                Qt.TransformationMode.SmoothTransformation
            )
            if isinstance(scaled, QPixmap) and not scaled.isNull():
                pixmap = scaled
                is_scaled = True
            else:
                is_scaled = True

        self._image_label._original_pixmap = original_pixmap
        self._image_label.set_pixmap(pixmap)
        self._image_label.set_fit_mode()

        # 显示缩放控件
        self._zoom_slider.setVisible(True)
        self._zoom_percent_label.setVisible(True)
        self._fit_view_btn.setVisible(True)
        self._view_original_btn.setVisible(is_scaled)

        # 更新缩放滑块和百分比
        self._update_zoom_controls(self._image_label.zoom_level())

        if is_scaled:
            self._file_info_label.setText(
                f"{file_item.size_display} · {original_pixmap.width()}×{original_pixmap.height()} · 已缩放预览"
            )
        else:
            self._file_info_label.setText(
                f"{file_item.size_display} · {pixmap.width()}×{pixmap.height()}"
            )

        self._content_stack.setCurrentWidget(self._image_preview)
        # 延迟一帧重新适配图片尺寸，确保 widget 已完成布局
        QTimer.singleShot(0, self._refit_current_image)

    def _show_gif_content(self, file_item):
        self._stop_gif()
        self._movie = QMovie(file_item.path)
        if not self._movie.isValid():
            self._show_unsupported(file_item, "无法加载 GIF 动画")
            return
        # GIF 模式：隐藏 _image_label，显示 _gif_label
        self._image_label.setVisible(False)
        self._gif_label.setVisible(True)
        self._gif_label.setMovie(self._movie)
        self._movie.start()
        self._view_original_btn.setVisible(False)
        self._zoom_slider.setVisible(False)
        self._zoom_percent_label.setVisible(False)
        self._fit_view_btn.setVisible(False)
        self._file_info_label.setText(f"{file_item.size_display} · GIF 动画")
        self._content_stack.setCurrentWidget(self._image_preview)

    def _stop_gif(self):
        if self._movie:
            self._movie.stop()
            self._movie = None

    def _fit_image_to_label(self, pixmap):
        """将图片设置为适应窗口模式"""
        if not isinstance(pixmap, QPixmap) or pixmap.isNull():
            return
        self._image_label.set_pixmap(pixmap)
        self._image_label.set_fit_mode()
        self._update_zoom_controls(self._image_label.zoom_level())

    def _refit_current_image(self):
        """延迟重新适配当前图片尺寸，用于首次预览时确保 widget 已完成布局"""
        if not self._current_file_item:
            return
        if self._content_stack.currentWidget() != self._image_preview:
            return
        if hasattr(self._image_label, '_original_pixmap') and not self._image_label._original_pixmap.isNull():
            self._image_label.set_pixmap(self._image_label._original_pixmap)
            self._image_label.set_fit_mode()
            self._update_zoom_controls(self._image_label.zoom_level())
        elif not self._image_label._pixmap.isNull():
            self._image_label.set_fit_mode()
            self._update_zoom_controls(self._image_label.zoom_level())

    def resizeEvent(self, event):
        super().resizeEvent(event)
        if self._movie and self._movie.state() == QMovie.MovieState.Running:
            return
        if (self._content_stack.currentWidget() == self._image_preview
                and self._current_file_item
                and self._current_file_item.extension.lower() in PREVIEWABLE_IMAGE_EXTS):
            if self._image_label.is_fit_mode():
                if hasattr(self._image_label, '_original_pixmap') and not self._image_label._original_pixmap.isNull():
                    self._image_label.set_pixmap(self._image_label._original_pixmap)
                self._image_label.set_fit_mode()
                self._update_zoom_controls(self._image_label.zoom_level())

    def _show_pdf_content(self, file_item, force=False):
        """显示 PDF 文件预览，以高 DPI 渲染一次，缩放时直接缩放图片无需重新渲染"""
        file_size = self._check_file_size(file_item, MAX_DOC_PREVIEW_SIZE_MB, force)
        if file_size is None and not force:
            return

        try:
            import fitz
        except ImportError:
            self._show_unsupported(file_item, "PDF 预览需要安装 PyMuPDF 库\n请运行: pip install PyMuPDF")
            return

        try:
            doc = fitz.open(file_item.path)
            self._pdf_total_pages = len(doc)
            doc.close()
        except Exception as e:
            logger.warning(f"PDF open failed: {file_item.path}, {type(e).__name__}")
            self._show_unsupported(file_item, "无法解析此 PDF 文件")
            return

        if self._pdf_total_pages == 0:
            self._show_unsupported(file_item, "PDF 文件为空")
            return

        self._pdf_pages_loaded = 0
        self._pdf_load_more_btn.setVisible(False)
        self._pdf_zoom = 1.0
        self._pdf_fit_mode = True
        self._pdf_original_pixmaps.clear()

        # 清空之前的页面标签
        for label in self._pdf_page_labels:
            self._pdf_pages_layout.removeWidget(label)
            label.deleteLater()
        self._pdf_page_labels.clear()

        info = f"{file_item.size_display} · 共 {self._pdf_total_pages} 页"
        self._file_info_label.setText(info)

        self._loading_text.setText("正在渲染 PDF...")
        self._loading_spinner.start()
        self._loading_progress.setVisible(True)
        self._content_stack.setCurrentWidget(self._loading_preview)

        # 以 150 DPI 渲染（足够清晰，缩放时用 Qt 平滑缩放）
        self._start_pdf_render(file_item.path, 0, DEFAULT_PDF_PAGES, 150)

    def _start_pdf_render(self, file_path, start_page, count, dpi=150):
        if self._pdf_render_worker and self._pdf_render_worker.isRunning():
            self._pdf_render_worker.cancel()
            self._pdf_render_worker.wait(2000)

        self._pdf_render_worker = _PdfRenderWorker(file_path, count, start_page, dpi)
        self._pdf_render_worker.page_rendered.connect(self._on_pdf_page_rendered)
        self._pdf_render_worker.render_finished.connect(self._on_pdf_render_finished)
        self._pdf_render_worker.start()

    def _on_pdf_page_rendered(self, page_idx, png_data, img_w, img_h):
        """收到渲染好的 PDF 页面，存储原始 pixmap 并按当前缩放显示"""
        if not self._current_file_item:
            return

        pixmap = QPixmap()
        pixmap.loadFromData(png_data, "PNG")
        if pixmap.isNull():
            return

        # 存储原始高 DPI pixmap
        self._pdf_original_pixmaps[page_idx] = pixmap

        # 按当前缩放级别显示
        self._display_pdf_page(page_idx, pixmap)

    def _display_pdf_page(self, page_idx: int, original_pixmap: QPixmap):
        """按当前缩放级别显示 PDF 页面"""
        # 计算目标显示宽度
        viewport_w = self._pdf_scroll_area.viewport().width()
        if viewport_w < 100:
            viewport_w = self._pdf_scroll_area.width()
        if viewport_w < 100:
            viewport_w = 400

        if self._pdf_fit_mode:
            # 适应窗口：宽度刚好等于视口宽度（减去边距）
            target_w = viewport_w - 8
        else:
            # 按缩放比例：基于适应窗口时的宽度乘以缩放因子
            fit_w = viewport_w - 8
            target_w = int(fit_w * self._pdf_zoom)

        if target_w < 50:
            target_w = 50

        scaled = original_pixmap.scaled(
            target_w, original_pixmap.height() * target_w // original_pixmap.width(),
            Qt.AspectRatioMode.KeepAspectRatio,
            Qt.TransformationMode.SmoothTransformation
        )

        # 查找或创建页面容器
        if page_idx < len(self._pdf_page_labels) and self._pdf_page_labels[page_idx] is not None:
            container = self._pdf_page_labels[page_idx]
            img_label = container.findChild(QLabel, "pdfPageImg")
            if img_label:
                img_label.setPixmap(scaled)
            return

        # 创建新的页面容器
        page_container = QWidget()
        page_layout = QVBoxLayout(page_container)
        page_layout.setContentsMargins(0, 0, 0, 4)
        page_layout.setSpacing(2)

        img_label = QLabel()
        img_label.setObjectName("pdfPageImg")
        img_label.setPixmap(scaled)
        img_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
        img_label.setStyleSheet(f"border:1px solid {COLORS.BORDER_DEFAULT}; border-radius:{RADIUS.SMALL}px; background:white;")

        page_num = QLabel(f"第 {page_idx + 1} 页")
        page_num.setAlignment(Qt.AlignmentFlag.AlignCenter)
        page_num.setStyleSheet(f"color: {COLORS.TEXT_TERTIARY}; font-size: {FONT.MICRO_PT - 1}px; border:none; background:transparent;")

        page_layout.addWidget(img_label)
        page_layout.addWidget(page_num)

        while len(self._pdf_page_labels) <= page_idx:
            self._pdf_page_labels.append(None)
        self._pdf_page_labels[page_idx] = page_container
        self._pdf_pages_layout.insertWidget(page_idx, page_container)
        self._pdf_pages_loaded = max(self._pdf_pages_loaded, page_idx + 1)

    def _on_pdf_render_finished(self, total_pages):
        self._loading_spinner.stop()
        self._loading_progress.setVisible(False)
        if self._current_file_item:
            self._content_stack.setCurrentWidget(self._pdf_preview)
        # 显示缩放控件
        self._pdf_fit_btn.setVisible(not self._pdf_fit_mode)
        self._pdf_zoom_label.setVisible(True)
        self._pdf_zoom_slider.setVisible(True)
        self._pdf_zoom_slider.blockSignals(True)
        self._pdf_zoom_slider.setValue(int(self._pdf_zoom * 100))
        self._pdf_zoom_slider.blockSignals(False)
        self._pdf_zoom_label.setText(f"{int(self._pdf_zoom * 100)}%")

        if self._pdf_pages_loaded < self._pdf_total_pages:
            self._pdf_load_more_btn.setVisible(True)
            self._pdf_load_more_btn.setEnabled(True)
            self._pdf_load_more_btn.setText(
                f"加载更多页面（剩余 {self._pdf_total_pages - self._pdf_pages_loaded} 页）"
            )
        else:
            self._pdf_load_more_btn.setVisible(False)

    def _on_pdf_load_more(self):
        if not self._current_file_item:
            return
        self._pdf_load_more_btn.setText("正在加载...")
        self._pdf_load_more_btn.setEnabled(False)
        self._start_pdf_render(self._current_file_item.path, self._pdf_pages_loaded, 10, 150)

    def _on_pdf_scroll_wheel(self, event):
        """PDF 滚轮事件：Ctrl+滚轮缩放，普通滚轮滚动"""
        if event.modifiers() & Qt.KeyboardModifier.ControlModifier:
            delta = event.angleDelta().y()
            if delta > 0:
                new_zoom = self._pdf_zoom * 1.15
            else:
                new_zoom = self._pdf_zoom / 1.15
            new_zoom = max(0.25, min(5.0, new_zoom))
            self._pdf_set_zoom(new_zoom)
            event.accept()
        else:
            QScrollArea.wheelEvent(self._pdf_scroll_area, event)

    def _pdf_set_zoom(self, zoom: float):
        """设置 PDF 缩放级别，直接缩放图片无需重新渲染"""
        if not self._current_file_item:
            return
        self._pdf_zoom = zoom
        self._pdf_fit_mode = False

        # 更新缩放控件
        self._pdf_zoom_slider.blockSignals(True)
        self._pdf_zoom_slider.setValue(int(zoom * 100))
        self._pdf_zoom_slider.blockSignals(False)
        self._pdf_zoom_label.setText(f"{int(zoom * 100)}%")
        self._pdf_fit_btn.setVisible(True)

        # 直接缩放所有已有页面，无需重新渲染
        self._pdf_refresh_display()

    def _pdf_refresh_display(self):
        """刷新所有 PDF 页面的显示（按当前缩放级别缩放原始 pixmap）"""
        for page_idx, original_pixmap in self._pdf_original_pixmaps.items():
            if page_idx < len(self._pdf_page_labels) and self._pdf_page_labels[page_idx] is not None:
                self._display_pdf_page(page_idx, original_pixmap)

    def _on_pdf_zoom_slider_changed(self, value: int):
        """PDF 缩放滑块值变化"""
        self._pdf_set_zoom(value / 100.0)

    def _on_pdf_fit_view(self):
        """PDF 适应窗口模式"""
        self._pdf_zoom = 1.0
        self._pdf_fit_mode = True
        self._pdf_fit_btn.setVisible(False)
        self._pdf_zoom_slider.blockSignals(True)
        self._pdf_zoom_slider.setValue(100)
        self._pdf_zoom_slider.blockSignals(False)
        self._pdf_zoom_label.setText("100%")
        self._pdf_refresh_display()

    def _show_excel_content(self, file_item, force=False):
        file_size = self._check_file_size(file_item, MAX_DOC_PREVIEW_SIZE_MB, force)
        if file_size is None and not force:
            return

        # 清空之前的 tab
        while self._excel_tab.count() > 0:
            self._excel_tab.removeTab(0)

        self._excel_text_edit.clear()
        self._excel_render_mode = False
        self._excel_mode_slider.set_current_index(0)
        self._excel_tab.setVisible(False)
        self._excel_text_edit.setVisible(True)

        try:
            from openpyxl import load_workbook
            wb = load_workbook(file_item.path, read_only=True, data_only=True)
            all_text = []
            total_chars = 0
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                all_text.append(f"=== {sheet_name} ===")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else '' for c in row]
                    line = '\t'.join(cells)
                    all_text.append(line)
                    total_chars += len(line)
                    if total_chars > MAX_DOC_CHARS and not force:
                        all_text.append("... 已截断")
                        break
                if total_chars > MAX_DOC_CHARS and not force:
                    break
            wb.close()
            self._excel_text_edit.setPlainText('\n'.join(all_text))
            cursor = self._excel_text_edit.textCursor()
            cursor.movePosition(QTextCursor.MoveOperation.Start)
            self._excel_text_edit.setTextCursor(cursor)
        except Exception as e:
            logger.warning(f"Excel text preview failed: {file_item.path}, {type(e).__name__}")
            self._excel_text_edit.setPlainText("无法读取文本内容")

        self._file_info_label.setText(f"{file_item.size_display} · Excel 文档")
        self._content_stack.setCurrentWidget(self._excel_preview)

    def _on_toggle_excel_mode(self, index: int = 1):
        if not self._current_file_item:
            return
        self._excel_render_mode = (index == 1)
        if self._excel_render_mode:
            self._excel_text_edit.setVisible(False)
            self._excel_tab.setVisible(True)

            self._loading_text.setText("正在加载 Excel...")
            self._loading_spinner.start()
            self._loading_progress.setVisible(True)
            self._content_stack.setCurrentWidget(self._loading_preview)

            if self._excel_render_worker and self._excel_render_worker.isRunning():
                self._excel_render_worker.cancel()
                self._excel_render_worker.wait(2000)

            # 清空之前的 tab
            while self._excel_tab.count() > 0:
                self._excel_tab.removeTab(0)

            self._excel_render_worker = _ExcelDataWorker(self._current_file_item.path, EXCEL_RENDER_PAGES)
            self._excel_render_worker.sheet_data_ready.connect(self._on_excel_sheet_data_ready)
            self._excel_render_worker.all_sheets_done.connect(self._on_excel_data_done)
            self._excel_render_worker.start()
        else:
            self._excel_text_edit.setVisible(True)
            self._excel_tab.setVisible(False)
            self._content_stack.setCurrentWidget(self._excel_preview)

    def _create_excel_table(self, headers, rows, sheet_name) -> QTableWidget:
        """创建带样式的 QTableWidget 用于 Excel 预览"""
        table = QTableWidget(len(rows), len(headers))
        table.setHorizontalHeaderLabels(headers)
        table.setEditTriggers(QTableWidget.EditTrigger.NoEditTriggers)
        table.setSelectionBehavior(QTableWidget.SelectionBehavior.SelectItems)
        table.setAlternatingRowColors(True)
        table.verticalHeader().setVisible(False)
        table.horizontalHeader().setStretchLastSection(True)
        table.setHorizontalScrollMode(QTableWidget.ScrollMode.ScrollPerPixel)
        table.setVerticalScrollMode(QTableWidget.ScrollMode.ScrollPerPixel)

        cell_font = QFont("Microsoft YaHei", FONT.MICRO_PT)
        header_font = QFont("Microsoft YaHei", FONT.MICRO_PT)
        header_font.setBold(True)

        for r, row_data in enumerate(rows):
            for c, cell_text in enumerate(row_data):
                item = QTableWidgetItem(cell_text)
                item.setFont(cell_font)
                item.setForeground(QColor(COLORS.TEXT_SECONDARY))
                item.setTextAlignment(Qt.AlignmentFlag.AlignVCenter | Qt.AlignmentFlag.AlignLeft)
                table.setItem(r, c, item)

        # 设置列头字体
        for c in range(len(headers)):
            header_item = table.horizontalHeaderItem(c)
            if header_item:
                header_item.setFont(header_font)
                header_item.setForeground(QColor(COLORS.TEXT_TERTIARY))

        # 自适应列宽
        table.resizeColumnsToContents()
        # 限制列宽范围
        for c in range(len(headers)):
            w = table.columnWidth(c)
            if w < 50:
                table.setColumnWidth(c, 50)
            elif w > 300:
                table.setColumnWidth(c, 300)

        table.setStyleSheet(f"""
            QTableWidget {{
                background-color: {COLORS.BG_PRIMARY};
                border: none;
                gridline-color: {COLORS.BORDER_DEFAULT};
                selection-background-color: {COLORS.BRAND_LIGHT_BG};
                selection-color: {COLORS.TEXT_PRIMARY};
                font-size: {FONT.MICRO_PT}px;
            }}
            QTableWidget::item {{
                padding: 4px 8px;
                border-bottom: 1px solid {COLORS.BORDER_DEFAULT};
                border-right: 1px solid {COLORS.BORDER_DEFAULT};
            }}
            QTableWidget::item:alternate {{
                background-color: {COLORS.BG_SECONDARY};
            }}
            QHeaderView::section {{
                background-color: {COLORS.BG_TERTIARY};
                color: {COLORS.TEXT_TERTIARY};
                border: none;
                border-bottom: 1px solid {COLORS.BORDER_DEFAULT};
                border-right: 1px solid {COLORS.BORDER_DEFAULT};
                padding: 4px 8px;
                font-weight: bold;
                font-size: {FONT.MICRO_PT}px;
            }}
            {scrollbar_style()}
        """)

        return table

    def _on_excel_sheet_data_ready(self, idx, headers, rows, sheet_name):
        if not self._current_file_item:
            return
        table = self._create_excel_table(headers, rows, sheet_name)
        self._excel_tab.addTab(table, sheet_name)

    def _on_excel_data_done(self, total_sheets):
        self._loading_spinner.stop()
        self._loading_progress.setVisible(False)
        if self._current_file_item:
            self._content_stack.setCurrentWidget(self._excel_preview)
        if total_sheets > EXCEL_RENDER_PAGES:
            self._file_info_label.setText(
                f"{self._current_file_item.size_display} · Excel 文档 · 显示前 {EXCEL_RENDER_PAGES} 个工作表（共 {total_sheets} 个）"
            )

    def _show_word_content(self, file_item, force=False):
        try:
            file_size = os.path.getsize(file_item.path)
        except OSError:
            self._show_unsupported(file_item, "无法读取文件")
            return

        if file_size > WORD_LARGE_FILE_CONFIRM_MB * 1024 * 1024 and not force:
            self._unsupported_icon.setPixmap(
                QIcon(f"icons/{FILE_ICON_MAP.get(file_item.extension.lower(), 'doctype/File.svg')}").pixmap(QSize(48, 48))
            )
            self._unsupported_text.setText("文件较大")
            self._unsupported_hint.setText(
                f"文件大小 {file_item.size_display}，预览可能需要较长时间\n是否仍要预览？"
            )
            self._force_preview_btn.setVisible(True)
            self._force_preview_btn.setText("预览文件")
            self._content_stack.setCurrentWidget(self._unsupported_preview)
            return

        file_size_val = self._check_file_size(file_item, MAX_DOC_PREVIEW_SIZE_MB, force)
        if file_size_val is None and not force:
            return

        try:
            from docx import Document
            doc = Document(file_item.path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            content = '\n'.join(paragraphs)

            if len(content) > MAX_DOC_CHARS and not force:
                content = content[:MAX_DOC_CHARS]
                content += f"\n\n... 已截断，共 {len(''.join(p.text for p in doc.paragraphs))} 字符"

            self._text_edit.setPlainText(content)
            self._show_more_btn.setVisible(False)
            self._md_mode_slider.setVisible(False)
            self._md_browser.setVisible(False)
            self._text_edit.setVisible(True)

            cursor = self._text_edit.textCursor()
            cursor.movePosition(QTextCursor.MoveOperation.Start)
            self._text_edit.setTextCursor(cursor)

            self._file_info_label.setText(f"{file_item.size_display} · Word 文档 · 纯文本预览")
            self._content_stack.setCurrentWidget(self._text_preview)
        except Exception as e:
            logger.warning(f"Word 预览失败: {file_item.path}, {type(e).__name__}")
            self._show_unsupported(file_item, "无法解析此 Word 文件\n可能是不受支持的格式")

    def _show_ppt_content(self, file_item, force=False):
        file_size = self._check_file_size(file_item, MAX_DOC_PREVIEW_SIZE_MB, force)
        if file_size is None and not force:
            return

        self._ppt_text_edit.clear()

        try:
            from pptx import Presentation
            prs = Presentation(file_item.path)
            all_text = []
            total_chars = 0
            slide_count = len(prs.slides)

            for i, slide in enumerate(prs.slides):
                all_text.append(f"=== 第 {i + 1} 页 ===")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                all_text.append(text)
                                total_chars += len(text)
                if total_chars > MAX_DOC_CHARS and not force:
                    all_text.append("... 已截断")
                    break

            content = '\n'.join(all_text)
            self._ppt_text_edit.setPlainText(content)
            cursor = self._ppt_text_edit.textCursor()
            cursor.movePosition(QTextCursor.MoveOperation.Start)
            self._ppt_text_edit.setTextCursor(cursor)

            self._file_info_label.setText(f"{file_item.size_display} · PPT 文档 · 共 {slide_count} 页")
        except Exception as e:
            logger.warning(f"PPT text preview failed: {file_item.path}, {type(e).__name__}")
            self._ppt_text_edit.setPlainText("无法读取文本内容")

        self._content_stack.setCurrentWidget(self._ppt_preview)

    def _show_video_content(self, file_item, force=False):
        """显示视频文件属性信息预览（类似右键属性对话框）"""
        file_size = self._check_file_size(file_item, MAX_MEDIA_PREVIEW_SIZE_MB, force)
        if file_size is None and not force:
            return
        self._show_media_properties(file_item, "视频文件")

    def _show_audio_content(self, file_item, force=False):
        """显示音频文件属性信息预览（类似右键属性对话框）"""
        file_size = self._check_file_size(file_item, MAX_MEDIA_PREVIEW_SIZE_MB, force)
        if file_size is None and not force:
            return
        self._show_media_properties(file_item, "音频文件")

    def _show_media_properties(self, file_item, type_label):
        """显示媒体文件属性信息，包含封面图和详细属性列表"""
        ext = file_item.extension.lower()
        is_video = ext in VIDEO_EXTENSIONS

        self._media_cover_label.setVisible(False)

        # 先显示基本信息表格
        basic_rows = [
            ("文件名", file_item.name),
            ("文件大小", file_item.size_display),
            ("文件路径", file_item.path),
            ("修改时间", file_item.modified_date),
        ]
        if is_video:
            basic_rows.insert(2, ("文件类型", "视频文件"))
        else:
            basic_rows.insert(2, ("文件类型", "音频文件"))
        self._populate_media_info_table(basic_rows)
        self._media_info_table.setVisible(True)

        self._open_system_player_btn.setVisible(True)
        self._media_fallback.setVisible(True)
        self._file_info_label.setText(f"{file_item.size_display} · {type_label}")
        self._content_stack.setCurrentWidget(self._media_preview)

        # 异步提取媒体信息和封面，避免阻塞 UI 线程
        from concurrent.futures import ThreadPoolExecutor

        def _fetch_info():
            media_info = self._get_media_info(file_item.path)
            thumbnail = self._get_video_thumbnail(file_item.path) if is_video else self._get_audio_cover(file_item.path)
            return media_info, thumbnail

        def _on_info_ready(future):
            try:
                media_info, thumbnail = future.result()
            except Exception:
                return
            if self._current_file_item != file_item:
                return

            # 构建详细属性列表：媒体信息在前，文件基本信息在后
            rows = []
            if media_info:
                if "duration" in media_info:
                    rows.append(("时长", media_info["duration"]))
                if is_video:
                    if "resolution" in media_info:
                        rows.append(("分辨率", media_info["resolution"]))
                    if "frame_width" in media_info:
                        rows.append(("帧宽度", media_info["frame_width"]))
                    if "frame_height" in media_info:
                        rows.append(("帧高度", media_info["frame_height"]))
                    if "display_aspect_ratio" in media_info:
                        rows.append(("宽高比", media_info["display_aspect_ratio"]))
                    if "fps" in media_info:
                        rows.append(("帧速率", media_info["fps"]))
                    if "scan_type" in media_info:
                        rows.append(("扫描方式", media_info["scan_type"]))
                    if "video_codec" in media_info:
                        codec_str = media_info["video_codec"]
                        if "video_profile" in media_info:
                            codec_str += f" ({media_info['video_profile']})"
                        rows.append(("视频编码", codec_str))
                    if "video_level" in media_info:
                        rows.append(("编码级别", media_info["video_level"]))
                    if "pixel_format" in media_info:
                        rows.append(("像素格式", media_info["pixel_format"]))
                    if "color_space" in media_info:
                        rows.append(("色彩空间", media_info["color_space"]))
                    if "color_range" in media_info:
                        rows.append(("色彩范围", media_info["color_range"]))
                    if "video_bitrate" in media_info:
                        rows.append(("视频比特率", media_info["video_bitrate"]))
                    if "video_stream_count" in media_info:
                        rows.append(("视频流数", media_info["video_stream_count"]))
                if "audio_codec" in media_info:
                    codec_str = media_info["audio_codec"]
                    if "audio_profile" in media_info:
                        codec_str += f" ({media_info['audio_profile']})"
                    rows.append(("音频编码", codec_str))
                if "audio_language" in media_info:
                    rows.append(("音频语言", media_info["audio_language"]))
                if "sample_rate" in media_info:
                    rows.append(("采样率", media_info["sample_rate"]))
                if "channels" in media_info:
                    rows.append(("声道", media_info["channels"]))
                if "audio_bitrate" in media_info:
                    rows.append(("音频比特率", media_info["audio_bitrate"]))
                if "bits_per_sample" in media_info:
                    rows.append(("位深度", media_info["bits_per_sample"]))
                if "audio_stream_count" in media_info:
                    rows.append(("音频流数", media_info["audio_stream_count"]))
                if "bitrate" in media_info:
                    rows.append(("总比特率", media_info["bitrate"]))
                if "container_format" in media_info:
                    rows.append(("容器格式", media_info["container_format"]))
                if "creation_time" in media_info:
                    rows.append(("创建时间", media_info["creation_time"]))
                # 媒体标签
                if "tag_title" in media_info:
                    rows.append(("标题", media_info["tag_title"]))
                if "tag_artist" in media_info:
                    rows.append(("艺术家", media_info["tag_artist"]))
                if "tag_album" in media_info:
                    rows.append(("专辑", media_info["tag_album"]))
            # 文件基本信息放在最后
            rows.append(("文件名", file_item.name))
            rows.append(("文件大小", file_item.size_display))
            rows.append(("文件路径", file_item.path))
            rows.append(("修改时间", file_item.modified_date))

            self._populate_media_info_table(rows)

            # 更新封面图
            if thumbnail and not thumbnail.isNull():
                max_w = self._media_cover_label.width()
                max_h = self._media_cover_label.height()
                if max_w < 10:
                    max_w = 280
                if max_h < 10:
                    max_h = 200
                scaled = thumbnail.scaled(
                    max_w, max_h,
                    Qt.AspectRatioMode.KeepAspectRatio,
                    Qt.TransformationMode.SmoothTransformation
                )
                self._media_cover_label.setPixmap(scaled)
                self._media_cover_label.setVisible(True)

        executor = ThreadPoolExecutor(max_workers=1)
        future = executor.submit(_fetch_info)
        future.add_done_callback(lambda f: QTimer.singleShot(0, lambda: _on_info_ready(f)))

    def _populate_media_info_table(self, rows: list):
        """填充媒体属性信息表格

        Args:
            rows: 属性键值对列表，如 [("文件名", "test.mp4"), ("大小", "10MB")]
        """
        self._media_info_table.setRowCount(len(rows))
        self._media_info_table.setColumnCount(2)
        self._media_info_table.setColumnWidth(0, 100)

        key_font = QFont("Microsoft YaHei", FONT.MICRO_PT - 1)
        key_font.setBold(True)
        value_font = QFont("Microsoft YaHei", FONT.MICRO_PT - 1)

        for i, (key, value) in enumerate(rows):
            key_item = QTableWidgetItem(key)
            key_item.setFont(key_font)
            key_item.setForeground(QColor(COLORS.TEXT_TERTIARY))
            key_item.setFlags(key_item.flags() & ~Qt.ItemFlag.ItemIsSelectable)

            value_item = QTableWidgetItem(str(value))
            value_item.setFont(value_font)
            value_item.setForeground(QColor(COLORS.TEXT_SECONDARY))
            value_item.setFlags(value_item.flags() & ~Qt.ItemFlag.ItemIsSelectable)

            self._media_info_table.setItem(i, 0, key_item)
            self._media_info_table.setItem(i, 1, value_item)

    def _get_audio_cover(self, file_path: str) -> QPixmap | None:
        """尝试从音频文件中提取封面图

        Args:
            file_path: 音频文件路径

        Returns:
            封面图 QPixmap，如果提取失败则返回 None
        """
        # 尝试使用 ffprobe 提取封面
        try:
            tmp_fd, tmp_path = tempfile.mkstemp(suffix=".jpg")
            os.close(tmp_fd)
            try:
                result = subprocess.run(
                    [
                        "ffmpeg", "-y",
                        "-i", file_path,
                        "-an", "-vcodec", "copy",
                        tmp_path
                    ],
                    capture_output=True, text=True, timeout=15,
                    creationflags=subprocess.CREATE_NO_WINDOW
                )
                if result.returncode == 0 and os.path.getsize(tmp_path) > 0:
                    pixmap = QPixmap(tmp_path)
                    if not pixmap.isNull():
                        return pixmap
            finally:
                try:
                    os.unlink(tmp_path)
                except OSError:
                    pass
        except (FileNotFoundError, OSError):
            logger.debug("ffmpeg 不可用，跳过音频封面提取")
        except subprocess.TimeoutExpired:
            logger.warning("ffmpeg 音频封面提取超时")
        except Exception as e:
            logger.warning(f"音频封面提取失败: {type(e).__name__}")
        return None

    def _show_archive_content(self, file_item):
        ext = file_item.extension.lower()
        entries = []
        dir_count = 0
        file_count = 0

        try:
            if ext == '.zip':
                with zipfile.ZipFile(file_item.path, 'r') as zf:
                    try:
                        zf.testzip()
                    except RuntimeError:
                        self._show_unsupported(file_item, "此压缩包有密码保护\n无法预览加密的压缩包")
                        return
                    for info in zf.infolist():
                        name = info.filename
                        if name.endswith('/'):
                            dir_count += 1
                            entries.append((name.rstrip('/'), True))
                        else:
                            file_count += 1
                            entries.append((os.path.basename(name), False))
            elif ext == '.rar':
                try:
                    import rarfile
                    with rarfile.RarFile(file_item.path, 'r') as rf:
                        for info in rf.infolist():
                            name = info.filename
                            if info.is_dir():
                                dir_count += 1
                                entries.append((name.rstrip('/'), True))
                            else:
                                file_count += 1
                                entries.append((os.path.basename(name), False))
                except ImportError:
                    self._show_unsupported(file_item, "RAR 预览需要安装 rarfile 库\n请运行: pip install rarfile\n可双击使用解压软件打开")
                    return
            elif ext == '.7z':
                try:
                    import py7zr
                    with py7zr.SevenZipFile(file_item.path, 'r') as sz:
                        for info in sz.list():
                            name = info.filename
                            if info.is_directory:
                                dir_count += 1
                                entries.append((name.rstrip('/'), True))
                            else:
                                file_count += 1
                                entries.append((os.path.basename(name), False))
                except ImportError:
                    self._show_unsupported(file_item, "7z 预览需要安装 py7zr 库\n请运行: pip install py7zr\n可双击使用解压软件打开")
                    return
            elif ext in {'.tar', '.gz', '.bz2', '.xz', '.tgz', '.tbz2'}:
                try:
                    with tarfile.open(file_item.path, 'r:*') as tf:
                        for member in tf.getmembers():
                            if member.isdir():
                                dir_count += 1
                                entries.append((os.path.basename(member.name) or member.name, True))
                            else:
                                file_count += 1
                                entries.append((os.path.basename(member.name), False))
                except tarfile.TarError:
                    self._show_unsupported(file_item, "无法解析此压缩包")
                    return
        except Exception as e:
            logger.warning(f"压缩包预览失败: {file_item.path}, {type(e).__name__}")
            self._show_unsupported(file_item, "无法读取压缩包内容")
            return

        self._show_archive_list(file_item, entries, dir_count, file_count)

    def _show_archive_list(self, file_item, entries, dir_count, file_count):
        self.title_label.setText(file_item.name)
        self.icon_label.setPixmap(
            QIcon(f"icons/{FILE_ICON_MAP.get(file_item.extension.lower(), 'doctype/Zip.svg')}").pixmap(QSize(20, 20))
        )

        while self._folder_list_layout.count():
            child = self._folder_list_layout.takeAt(0)
            if child.widget():
                child.widget().deleteLater()

        entries.sort(key=lambda e: (not e[1], e[0].lower()))

        stats_parts = []
        if dir_count > 0:
            stats_parts.append(f"{dir_count} 个文件夹")
        if file_count > 0:
            stats_parts.append(f"{file_count} 个文件")
        stats_text = " · ".join(stats_parts) if stats_parts else "空压缩包"
        self._folder_stats_label.setText(f"压缩包内容 · {stats_text}")
        self._folder_stats_label.setStyleSheet(label_caption_style())

        max_entries = 50
        display_entries = entries[:max_entries]
        for name, is_dir in display_entries:
            row = _FolderEntryRow(name, is_dir)
            self._folder_list_layout.addWidget(row)

        if len(entries) > max_entries:
            more_label = QLabel(f"  ... 还有 {len(entries) - max_entries} 项未显示")
            more_label.setStyleSheet(f"""
                color: {COLORS.TEXT_PLACEHOLDER}; font-size: {FONT.MICRO_PT}px;
                border: none; background: transparent; padding: 4px 6px;
            """)
            self._folder_list_layout.addWidget(more_label)

        if not entries:
            empty_label = QLabel("压缩包为空")
            empty_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
            empty_label.setStyleSheet(f"""
                color: {COLORS.TEXT_PLACEHOLDER}; font-size: {FONT.CAPTION_PT}px;
                border: none; background: transparent; padding: 16px;
            """)
            self._folder_list_layout.addWidget(empty_label)

        self._file_info_label.setText(file_item.size_display)
        self._content_stack.setCurrentWidget(self._folder_preview)

    def _show_unsupported(self, file_item, hint_text=""):
        self._unsupported_icon.setPixmap(
            QIcon(f"icons/{FILE_ICON_MAP.get(file_item.extension.lower(), 'doctype/File.svg')}").pixmap(QSize(48, 48))
        )
        self._unsupported_hint.setText(hint_text)
        self._force_preview_btn.setVisible(False)
        self._content_stack.setCurrentWidget(self._unsupported_preview)

    def _show_folder_preview(self, file_item):
        self._hide_all_content()
        self.title_label.setText(file_item.name)
        self._file_info_label.setText(file_item.item_count_display)
        self.icon_label.setPixmap(QIcon("icons/doctype/Folder.svg").pixmap(QSize(20, 20)))

        while self._folder_list_layout.count():
            child = self._folder_list_layout.takeAt(0)
            if child.widget():
                child.widget().deleteLater()

        dir_path = file_item.path
        dir_count = 0
        file_count = 0
        entries = []

        try:
            for entry in os.scandir(dir_path):
                try:
                    if entry.is_dir(follow_symlinks=False):
                        dir_count += 1
                        entries.append((entry.name, True))
                    elif entry.is_file(follow_symlinks=False):
                        file_count += 1
                        entries.append((entry.name, False))
                except (PermissionError, OSError):
                    continue
        except (PermissionError, OSError):
            self._folder_stats_label.setText("无法访问此目录")
            self._folder_stats_label.setStyleSheet(f"color: {COLORS.TEXT_TERTIARY}; font-size: {FONT.CAPTION_PT}px; border: none; background: transparent;")
            self._content_stack.setCurrentWidget(self._folder_preview)
            return

        entries.sort(key=lambda e: (not e[1], e[0].lower()))

        stats_parts = []
        if dir_count > 0:
            stats_parts.append(f"{dir_count} 个文件夹")
        if file_count > 0:
            stats_parts.append(f"{file_count} 个文件")
        stats_text = " · ".join(stats_parts) if stats_parts else "空文件夹"
        self._folder_stats_label.setText(stats_text)
        self._folder_stats_label.setStyleSheet(label_caption_style())

        max_entries = 50
        display_entries = entries[:max_entries]
        for name, is_dir in display_entries:
            row = _FolderEntryRow(name, is_dir)
            self._folder_list_layout.addWidget(row)

        if len(entries) > max_entries:
            more_label = QLabel(f"  ... 还有 {len(entries) - max_entries} 项未显示")
            more_label.setStyleSheet(f"""
                color: {COLORS.TEXT_PLACEHOLDER}; font-size: {FONT.MICRO_PT}px;
                border: none; background: transparent; padding: 4px 6px;
            """)
            self._folder_list_layout.addWidget(more_label)

        if not entries:
            empty_label = QLabel("此文件夹为空")
            empty_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
            empty_label.setStyleSheet(f"""
                color: {COLORS.TEXT_PLACEHOLDER}; font-size: {FONT.CAPTION_PT}px;
                border: none; background: transparent; padding: 16px;
            """)
            self._folder_list_layout.addWidget(empty_label)

        self._content_stack.setCurrentWidget(self._folder_preview)

    def _on_show_more(self):
        if self._full_content:
            self._text_edit.setPlainText(self._full_content)
            self._showing_truncated = False
            self._show_more_btn.setVisible(False)
            ext = self._current_file_item.extension.lower() if self._current_file_item else ''
            self._highlighter = _create_highlighter(self._text_edit.document(), ext)
            self._file_info_label.setText(
                f"{self._current_file_item.size_display} · 显示全部内容"
            )

    def _on_toggle_md_mode(self, index: int = 1):
        self._md_source_mode = (index == 1)

        if self._full_content:
            lines = self._full_content.splitlines()
            truncated = len(lines) > MAX_PREVIEW_LINES
            self._render_markdown(self._full_content, truncated)

            if self._md_source_mode:
                self._file_info_label.setText(
                    f"{self._current_file_item.size_display} · 源码模式"
                )
            else:
                self._file_info_label.setText(
                    f"{self._current_file_item.size_display} · 预览模式"
                )

    def _show_html_content(self, file_item, force=False):
        """显示 HTML 文件预览，支持渲染/源码模式切换"""
        file_size = self._check_file_size(file_item, MAX_TEXT_PREVIEW_SIZE_MB, force)
        if file_size is None and not force:
            return

        from utils.encoding import read_text_file
        max_mb = MAX_TEXT_PREVIEW_SIZE_MB if not force else 50
        content = read_text_file(file_item.path, max_size_mb=max_mb)
        if content is None:
            self._show_unsupported(file_item, "无法读取文件内容")
            return

        self._full_content = content
        lines = content.splitlines()
        truncated = len(lines) > MAX_PREVIEW_LINES
        self._showing_truncated = truncated

        self._html_source_mode = False
        self._html_mode_slider.setVisible(True)
        self._html_mode_slider.set_current_index(0)
        self._md_mode_slider.setVisible(False)
        self._show_more_btn.setVisible(truncated)

        self._render_html(content, truncated)

        if truncated:
            self._file_info_label.setText(
                f"{file_item.size_display} · 渲染模式 · 显示前 {MAX_PREVIEW_LINES} 行"
            )
        else:
            self._file_info_label.setText(f"{file_item.size_display} · 渲染模式")

        self._content_stack.setCurrentWidget(self._text_preview)

    def _render_html(self, content, truncated):
        """渲染 HTML 内容，根据模式显示渲染结果或源码"""
        lines = content.splitlines()
        display_lines = lines[:MAX_PREVIEW_LINES] if truncated else lines
        display_content = '\n'.join(display_lines)

        if self._html_source_mode:
            self._text_edit.setPlainText(display_content)
            self._text_edit.setVisible(True)
            self._md_browser.setVisible(False)
            self._highlighter = _create_highlighter(self._text_edit.document(), '.html')
        else:
            # 渲染模式：将相对路径的图片转为绝对 file:// URL
            if self._current_file_item:
                file_dir = os.path.dirname(os.path.abspath(self._current_file_item.path))
                def _replace_src(match):
                    src = match.group(1)
                    if src and not src.startswith(('http://', 'https://', 'file://', 'data:', '#', 'mailto:')):
                        abs_path = os.path.normpath(os.path.join(file_dir, src))
                        file_url = 'file:///' + abs_path.replace('\\', '/')
                        return f'src="{file_url}"'
                    return match.group(0)
                display_content = re.sub(r'src="([^"]*)"', _replace_src, display_content)
                # 同样处理 href 中的相对路径
                def _replace_href(match):
                    href = match.group(1)
                    if href and not href.startswith(('http://', 'https://', 'file://', '#', 'mailto:', 'javascript:')):
                        abs_path = os.path.normpath(os.path.join(file_dir, href))
                        file_url = 'file:///' + abs_path.replace('\\', '/')
                        return f'href="{file_url}"'
                    return match.group(0)
                display_content = re.sub(r'href="([^"]*)"', _replace_href, display_content)

            styled_html = f"""
            <style>
                body {{ font-family: "Microsoft YaHei", sans-serif; color: {COLORS.TEXT_SECONDARY};
                       font-size: {FONT.CAPTION_PT}px; line-height: 1.6;
                       margin: 0; padding: 0; }}
                img {{ max-width: 100%; }}
                a {{ color: {COLORS.BRAND}; text-decoration: none; }}
                table {{ border-collapse: collapse; margin: 8px 0; }}
                th, td {{ border: 1px solid {COLORS.BORDER_DEFAULT}; padding: 6px 12px; }}
                th {{ background-color: {COLORS.BG_TERTIARY}; }}
            </style>
            {display_content}
            """
            self._md_browser.setHtml(styled_html)
            self._md_browser.setVisible(True)
            self._text_edit.setVisible(False)

    def _on_toggle_html_mode(self, index: int = 0):
        """切换 HTML 预览的渲染/源码模式"""
        self._html_source_mode = (index == 1)

        if self._full_content:
            lines = self._full_content.splitlines()
            truncated = len(lines) > MAX_PREVIEW_LINES
            self._render_html(self._full_content, truncated)

            if self._html_source_mode:
                self._file_info_label.setText(
                    f"{self._current_file_item.size_display} · 源码模式"
                )
            else:
                self._file_info_label.setText(
                    f"{self._current_file_item.size_display} · 渲染模式"
                )

    def _on_view_original_image(self):
        if hasattr(self._image_label, '_original_pixmap'):
            pm = self._image_label._original_pixmap
            if not isinstance(pm, QPixmap) or pm.isNull():
                return
            self._image_label.set_pixmap(pm)
            self._image_label.set_zoom(1.0)
            self._view_original_btn.setVisible(False)
            self._file_info_label.setText(
                f"{self._current_file_item.size_display} · {pm.width()}×{pm.height()} · 原始尺寸"
            )

    def _on_image_zoom_changed(self, zoom_level: float):
        """图片缩放级别变化时更新 UI"""
        percent = round(zoom_level * 100)
        self._zoom_percent_label.setText(f"{percent}%")
        # 阻断信号避免循环
        self._zoom_slider.blockSignals(True)
        self._zoom_slider.setValue(max(10, min(1000, percent)))
        self._zoom_slider.blockSignals(False)
        self._fit_view_btn.setVisible(not self._image_label.is_fit_mode())

    def _on_fit_view(self):
        """切换为适应窗口模式"""
        self._image_label.set_fit_mode()

    def _on_zoom_slider_changed(self, value: int):
        """缩放滑块值变化时设置缩放级别"""
        self._image_label.set_zoom(value / 100.0)

    def _update_zoom_controls(self, zoom_level: float):
        """更新缩放控件显示"""
        percent = round(zoom_level * 100)
        self._zoom_percent_label.setText(f"{percent}%")
        self._zoom_slider.blockSignals(True)
        self._zoom_slider.setValue(max(10, min(1000, percent)))
        self._zoom_slider.blockSignals(False)
        self._zoom_slider.setVisible(True)
        self._zoom_percent_label.setVisible(True)
        self._fit_view_btn.setVisible(not self._image_label.is_fit_mode())

    def _on_force_preview(self):
        if self._current_file_item:
            self._show_file_preview(self._current_file_item, force=True)

    def _do_delayed_preview(self):
        if self._pending_file_item:
            self._show_file_preview(self._pending_file_item, self._pending_force)
            self._pending_file_item = None

    def _cancel_workers(self):
        if self._pdf_render_worker and self._pdf_render_worker.isRunning():
            self._pdf_render_worker.cancel()
            self._pdf_render_worker.wait(2000)
            self._pdf_render_worker = None
        if self._excel_render_worker and self._excel_render_worker.isRunning():
            self._excel_render_worker.cancel()
            self._excel_render_worker.wait(2000)
            self._excel_render_worker = None

    def _stop_media(self):
        """停止媒体播放（保留接口，当前不再使用内联播放器）"""
        pass

    def _on_open_with_system_player(self):
        """使用系统默认播放器打开当前媒体文件"""
        if self._current_file_item and os.path.isfile(self._current_file_item.path):
            os.startfile(self._current_file_item.path)

    def _format_duration(self, seconds: float) -> str:
        """将秒数格式化为 HH:MM:SS 或 MM:SS 的时长字符串"""
        try:
            total = int(float(seconds))
        except (ValueError, TypeError):
            return "未知"
        hours = total // 3600
        minutes = (total % 3600) // 60
        secs = total % 60
        if hours > 0:
            return f"{hours:02d}:{minutes:02d}:{secs:02d}"
        return f"{minutes:02d}:{secs:02d}"

    def _format_bitrate(self, bps: str) -> str:
        """将比特率格式化为人类可读的字符串"""
        try:
            value = int(bps)
        except (ValueError, TypeError):
            return "未知"
        if value >= 1_000_000:
            return f"{value / 1_000_000:.1f} Mbps"
        elif value >= 1_000:
            return f"{value / 1_000:.0f} Kbps"
        return f"{value} bps"

    def _format_sample_rate(self, hz: str) -> str:
        """将采样率格式化为人类可读的字符串"""
        try:
            value = int(hz)
        except (ValueError, TypeError):
            return "未知"
        if value >= 1000:
            return f"{value / 1000:.1f} kHz"
        return f"{value} Hz"

    def _get_media_info(self, file_path: str) -> dict | None:
        """
        使用 ffprobe 提取媒体文件信息。

        Args:
            file_path: 媒体文件路径

        Returns:
            包含媒体元数据的字典，如果 ffprobe 不可用则返回 None
        """
        try:
            result = subprocess.run(
                [
                    "ffprobe", "-v", "quiet",
                    "-print_format", "json",
                    "-show_format", "-show_streams",
                    file_path
                ],
                capture_output=True, text=True, timeout=10,
                creationflags=subprocess.CREATE_NO_WINDOW
            )
            if result.returncode != 0:
                return None
            data = json.loads(result.stdout)
        except (FileNotFoundError, OSError):
            logger.debug("ffprobe 不可用，跳过媒体信息提取")
            return None
        except subprocess.TimeoutExpired:
            logger.warning("ffprobe 执行超时")
            return None
        except (json.JSONDecodeError, Exception) as e:
            logger.warning(f"ffprobe 输出解析失败: {type(e).__name__}")
            return None

        info = {}
        streams = data.get("streams", [])
        format_info = data.get("format", {})

        # 从 format 中提取通用信息
        format_name = format_info.get("format_name")
        if format_name:
            # 取第一个格式名（如 "mov,mp4,m4a,3gp,3g2,mj2" → "mov,mp4"）
            parts = format_name.split(",")
            info["container_format"] = parts[0].upper() if parts else format_name.upper()
        duration = format_info.get("duration")
        if duration:
            info["duration"] = self._format_duration(duration)
        bitrate = format_info.get("bit_rate")
        if bitrate:
            info["bitrate"] = self._format_bitrate(bitrate)

        # 从 format.tags 中提取标签信息
        tags = format_info.get("tags", {})
        if tags:
            title = tags.get("title")
            if title:
                info["tag_title"] = title
            artist = tags.get("artist")
            if artist:
                info["tag_artist"] = artist
            album = tags.get("album")
            if album:
                info["tag_album"] = album
            creation_time = tags.get("creation_time")
            if creation_time:
                info["creation_time"] = creation_time

        # 从视频流中提取信息
        video_count = 0
        audio_count = 0
        for stream in streams:
            if stream.get("codec_type") == "video":
                video_count += 1
                if "resolution" not in info:
                    width = stream.get("width")
                    height = stream.get("height")
                    if width and height:
                        info["resolution"] = f"{width}×{height}"
                        info["frame_width"] = f"{width} px"
                        info["frame_height"] = f"{height} px"
                    codec = stream.get("codec_name")
                    if codec:
                        codec_names = {
                            "h264": "H.264", "h265": "H.265", "hevc": "H.265/HEVC",
                            "vp8": "VP8", "vp9": "VP9", "av1": "AV1",
                            "mpeg4": "MPEG-4", "mpeg2video": "MPEG-2",
                            "wmv3": "WMV3", "msmpeg4v3": "MS MPEG-4 v3",
                        }
                        info["video_codec"] = codec_names.get(codec.lower(), codec.upper())
                    fps = stream.get("r_frame_rate")
                    if fps and "/" in fps:
                        try:
                            num, den = fps.split("/")
                            if int(den) != 0:
                                info["fps"] = f"{int(num) / int(den):.2f} fps"
                        except (ValueError, ZeroDivisionError):
                            pass
                    pixel_fmt = stream.get("pix_fmt")
                    if pixel_fmt:
                        info["pixel_format"] = pixel_fmt
                    video_br = stream.get("bit_rate")
                    if video_br:
                        info["video_bitrate"] = self._format_bitrate(video_br)
                    color_space = stream.get("color_space")
                    if color_space:
                        color_space_names = {
                            "bt709": "BT.709", "bt470bg": "BT.470BG",
                            "bt601": "BT.601", "smpte170m": "SMPTE 170M",
                            "smpte240m": "SMPTE 240M", "bt2020nc": "BT.2020NC",
                            "bt2020c": "BT.2020C",
                        }
                        info["color_space"] = color_space_names.get(color_space.lower(), color_space.upper())
                    color_range = stream.get("color_range")
                    if color_range:
                        color_range_names = {"tv": "Limited (TV)", "pc": "Full (PC)"}
                        info["color_range"] = color_range_names.get(color_range.lower(), color_range)
                    dar = stream.get("display_aspect_ratio")
                    if dar:
                        info["display_aspect_ratio"] = dar
                    profile = stream.get("profile")
                    if profile:
                        info["video_profile"] = profile
                    level = stream.get("level")
                    if level:
                        info["video_level"] = str(level)
                    # 扫描方式（逐行/隔行）
                    field_order = stream.get("field_order")
                    if field_order:
                        scan_names = {
                            "progressive": "逐行扫描",
                            "tt": "隔行扫描 (顶场优先)",
                            "bb": "隔行扫描 (底场优先)",
                            "tb": "隔行扫描 (顶场优先)",
                            "bt": "隔行扫描 (底场优先)",
                        }
                        scan_label = scan_names.get(field_order.lower())
                        if scan_label:
                            info["scan_type"] = scan_label

            elif stream.get("codec_type") == "audio":
                audio_count += 1
                if "audio_codec" not in info:
                    codec = stream.get("codec_name")
                    if codec:
                        codec_names = {
                            "aac": "AAC", "mp3": "MP3", "flac": "FLAC",
                            "pcm_s16le": "PCM", "vorbis": "Vorbis",
                            "opus": "Opus", "wmav2": "WMA",
                        }
                        info["audio_codec"] = codec_names.get(codec.lower(), codec.upper())
                    sample_rate = stream.get("sample_rate")
                    if sample_rate:
                        info["sample_rate"] = self._format_sample_rate(sample_rate)
                    channels = stream.get("channels")
                    if channels:
                        channel_map = {1: "单声道", 2: "立体声", 6: "5.1", 8: "7.1"}
                        info["channels"] = channel_map.get(channels, f"{channels} 声道")
                    audio_br = stream.get("bit_rate")
                    if audio_br:
                        info["audio_bitrate"] = self._format_bitrate(audio_br)
                    bps = stream.get("bits_per_sample")
                    if bps and int(bps) > 0:
                        info["bits_per_sample"] = f"{bps} bit"
                    audio_tags = stream.get("tags", {})
                    language = audio_tags.get("language")
                    if language:
                        lang_names = {
                            "chi": "中文", "zho": "中文", "zh": "中文",
                            "eng": "英语", "en": "英语", "jpn": "日语", "ja": "日语",
                            "kor": "韩语", "ko": "韩语", "fre": "法语", "fra": "法语",
                            "ger": "德语", "deu": "德语", "spa": "西班牙语",
                            "rus": "俄语", "und": "未知",
                        }
                        info["audio_language"] = lang_names.get(language.lower(), language.upper())
                    audio_profile = stream.get("profile")
                    if audio_profile:
                        info["audio_profile"] = audio_profile

        # 流数量统计
        if video_count > 0:
            info["video_stream_count"] = f"{video_count}"
        if audio_count > 0:
            info["audio_stream_count"] = f"{audio_count}"

        return info if info else None

    def _get_video_thumbnail(self, file_path: str) -> QPixmap | None:
        """
        使用 ffmpeg 提取视频缩略图。

        Args:
            file_path: 视频文件路径

        Returns:
            缩略图 QPixmap，如果提取失败则返回 None
        """
        try:
            tmp_fd, tmp_path = tempfile.mkstemp(suffix=".jpg")
            os.close(tmp_fd)
            try:
                result = subprocess.run(
                    [
                        "ffmpeg", "-y",
                        "-ss", "00:00:01",
                        "-i", file_path,
                        "-vframes", "1",
                        "-q:v", "2",
                        "-vf", "scale=320:-1",
                        tmp_path
                    ],
                    capture_output=True, text=True, timeout=15,
                    creationflags=subprocess.CREATE_NO_WINDOW
                )
                if result.returncode == 0 and os.path.getsize(tmp_path) > 0:
                    pixmap = QPixmap(tmp_path)
                    if not pixmap.isNull():
                        return pixmap
            finally:
                try:
                    os.unlink(tmp_path)
                except OSError:
                    pass
        except (FileNotFoundError, OSError):
            logger.debug("ffmpeg 不可用，跳过视频缩略图提取")
        except subprocess.TimeoutExpired:
            logger.warning("ffmpeg 缩略图提取超时")
        except Exception as e:
            logger.warning(f"视频缩略图提取失败: {type(e).__name__}")
        return None

    def clear_preview(self):
        self._cancel_workers()
        self._stop_media()
        self._stop_gif()
        self._loading_spinner.stop()
        self._loading_progress.setVisible(False)
        self._current_result = None
        self._full_content = None
        self._showing_truncated = False
        self._current_file_item = None
        self._highlighter = None
        self._pdf_doc = None
        self._preview_active = False
        self._show_empty()
