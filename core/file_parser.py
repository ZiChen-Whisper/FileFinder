import logging
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Optional, List, Set
from utils.encoding import read_text_file

logger = logging.getLogger(__name__)


class FileParser(ABC):
    """文件解析器基类"""

    @abstractmethod
    def can_parse(self, file_path: str) -> bool:
        """判断是否能够解析此文件"""
        pass

    @abstractmethod
    def parse(self, file_path: str) -> Optional[str]:
        """解析文件并返回文本内容"""
        pass

    @property
    def supported_extensions(self) -> Set[str]:
        """返回该解析器支持的文件扩展名集合"""
        return set()


class TextParser(FileParser):
    """纯文本解析器"""

    def __init__(self):
        self._text_exts = {'.txt', '.md', '.log', '.csv', '.json', '.xml',
                           '.yaml', '.yml', '.ini', '.cfg', '.conf', '.toml', '.env', '.gitignore',
                           '.py', '.js', '.ts', '.html', '.css', '.java',
                           '.c', '.cpp', '.h', '.go', '.rs', '.rb', '.php',
                           '.sh', '.bat', '.ps1', '.sql'}

    def can_parse(self, file_path: str) -> bool:
        """判断文件是否为纯文本文件"""
        return Path(file_path).suffix.lower() in self._text_exts

    def parse(self, file_path: str) -> Optional[str]:
        """解析纯文本文件内容"""
        return read_text_file(file_path)

    @property
    def supported_extensions(self) -> Set[str]:
        """返回支持的文本文件扩展名集合"""
        return self._text_exts


class PDFParser(FileParser):
    """PDF 文件解析器，提取内嵌文字内容"""

    def __init__(self):
        self._pdf_exts = {'.pdf'}

    def can_parse(self, file_path: str) -> bool:
        """判断文件是否为 PDF 文件"""
        return Path(file_path).suffix.lower() in self._pdf_exts

    def parse(self, file_path: str) -> Optional[str]:
        """解析 PDF 文件，提取全部页面的文字内容。

        Returns:
            各页文本用换行符连接的字符串，每页之间用空行分隔；
            如果 PDF 无文字内容或解析失败则返回 None。
        """
        try:
            import fitz
        except ImportError:
            logger.warning("PDF 解析需要 PyMuPDF 库，请运行: pip install PyMuPDF")
            return None

        try:
            doc = fitz.open(file_path)
            texts = []
            for page in doc:
                page_text = page.get_text("text")
                if page_text.strip():
                    texts.append(page_text)
            doc.close()

            if not texts:
                return None

            return "\n\n".join(texts)
        except Exception as e:
            logger.warning(f"PDF 解析失败: {file_path}, {type(e).__name__}")
            return None

    @property
    def supported_extensions(self) -> Set[str]:
        """返回支持的 PDF 文件扩展名集合"""
        return self._pdf_exts


class DocxParser(FileParser):
    """Word 文档解析器，提取 .docx 文件中的文字内容"""

    def __init__(self):
        self._docx_exts = {'.docx'}

    def can_parse(self, file_path: str) -> bool:
        """判断文件是否为 Word 文档"""
        return Path(file_path).suffix.lower() in self._docx_exts

    def parse(self, file_path: str) -> Optional[str]:
        """解析 Word 文档，提取全部段落文字内容。

        Returns:
            各段落文本用换行符连接的字符串；
            如果文档无文字内容或解析失败则返回 None。
        """
        try:
            from docx import Document
        except ImportError:
            logger.warning("Word 解析需要 python-docx 库，请运行: pip install python-docx")
            return None

        try:
            doc = Document(file_path)
            texts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    texts.append(para.text)

            # 提取表格中的文字
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        texts.append(" | ".join(row_text))

            if not texts:
                return None

            return "\n".join(texts)
        except Exception as e:
            logger.warning(f"Word 文档解析失败: {file_path}, {type(e).__name__}")
            return None

    @property
    def supported_extensions(self) -> Set[str]:
        """返回支持的 Word 文件扩展名集合"""
        return self._docx_exts


class XlsxParser(FileParser):
    """Excel 文件解析器，提取 .xlsx 文件中的文字内容"""

    def __init__(self):
        self._xlsx_exts = {'.xlsx'}

    def can_parse(self, file_path: str) -> bool:
        """判断文件是否为 Excel 文件"""
        return Path(file_path).suffix.lower() in self._xlsx_exts

    def parse(self, file_path: str) -> Optional[str]:
        """解析 Excel 文件，提取所有工作表中的文字内容。

        使用 read_only 模式降低内存占用。

        Returns:
            各单元格文本用制表符连接的字符串；
            如果文件无文字内容或解析失败则返回 None。
        """
        try:
            from openpyxl import load_workbook
        except ImportError:
            logger.warning("Excel 解析需要 openpyxl 库，请运行: pip install openpyxl")
            return None

        try:
            wb = load_workbook(file_path, read_only=True, data_only=True)
            texts = []
            sheet_count = len(wb.sheetnames)

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                texts.append(f"--- {sheet_name} ---")
                for row in ws.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None and str(cell).strip():
                            row_text.append(str(cell).strip())
                    if row_text:
                        texts.append("\t".join(row_text))

            wb.close()

            if len(texts) <= sheet_count:
                # 只有表头没有实际内容
                return None

            return "\n".join(texts)
        except Exception as e:
            logger.warning(f"Excel 文件解析失败: {file_path}, {type(e).__name__}")
            return None

    @property
    def supported_extensions(self) -> Set[str]:
        """返回支持的 Excel 文件扩展名集合"""
        return self._xlsx_exts


class PptxParser(FileParser):
    """PPT 文件解析器，提取 .pptx 文件中的文字内容"""

    def __init__(self):
        self._pptx_exts = {'.pptx'}

    def can_parse(self, file_path: str) -> bool:
        """判断文件是否为 PPT 文件"""
        return Path(file_path).suffix.lower() in self._pptx_exts

    def parse(self, file_path: str) -> Optional[str]:
        """解析 PPT 文件，提取所有幻灯片中的文字内容。

        Returns:
            各幻灯片文本用换行符连接的字符串；
            如果文件无文字内容或解析失败则返回 None。
        """
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("PPT 解析需要 python-pptx 库，请运行: pip install python-pptx")
            return None

        try:
            prs = Presentation(file_path)
            texts = []

            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                slide_texts.append(para.text.strip())
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_texts.append(" | ".join(row_text))

                if slide_texts:
                    texts.append(f"--- 幻灯片 {slide_idx} ---")
                    texts.extend(slide_texts)

            if not texts:
                return None

            return "\n".join(texts)
        except Exception as e:
            logger.warning(f"PPT 文件解析失败: {file_path}, {type(e).__name__}")
            return None

    @property
    def supported_extensions(self) -> Set[str]:
        """返回支持的 PPT 文件扩展名集合"""
        return self._pptx_exts


class ParserRegistry:
    """文件解析器注册表"""

    def __init__(self):
        self._parsers: List[FileParser] = [
            TextParser(),
            PDFParser(),
            DocxParser(),
            XlsxParser(),
            PptxParser(),
        ]

    def parse(self, file_path: str) -> Optional[str]:
        """根据文件类型选择合适的解析器"""
        # 跳过 Office 临时文件（以 ~$ 开头）
        filename = Path(file_path).name
        if filename.startswith('~$'):
            return None
        for parser in self._parsers:
            if parser.can_parse(file_path):
                return parser.parse(file_path)
        return None

    def can_parse(self, file_path: str) -> bool:
        """判断是否有解析器能处理此文件。

        Args:
            file_path: 文件路径

        Returns:
            是否可解析
        """
        filename = Path(file_path).name
        if filename.startswith('~$'):
            return False
        for parser in self._parsers:
            if parser.can_parse(file_path):
                return True
        return False

    def get_all_supported_extensions(self) -> Set[str]:
        """汇总所有解析器支持的文件扩展名集合"""
        exts = set()
        for parser in self._parsers:
            exts.update(parser.supported_extensions)
        return exts
